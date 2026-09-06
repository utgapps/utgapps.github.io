"""AI101 course builder.

Everything in this folder is generated. Never hand-edit the HTML - edit
course.py and re-run this. Same contract as camp-coding-projects/workbooks.py.

    python build.py

Emits, all derived from the one WEEKS structure in course.py:

    index.html          course hub
    week-01..15.html    the code every student should have by the end of week N
    teacher.html        full teaching curriculum, minute by minute
    workbook.html       15-chapter printable student homework book
    slides/week-NN.pptx one deck per week (needs `pip install python-pptx`)

Milestones are REPLAYED from the week ops rather than stored, so week N is
provably week N-1 plus that week's changes - they cannot drift apart, and
neither can the teacher guide, the homework or the slides.
"""

import difflib
import hashlib
import html
import json
import os
import re
import sys

import course

HERE = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(HERE, "slides")

LOGO = "https://s3.us-west-1.amazonaws.com/utg.pictures.videos/UTGWeb/utglogoh.svg"
FONT = '<link rel="preconnect" href="https://fonts.googleapis.com"><link rel="preconnect" href="https://fonts.gstatic.com" crossorigin><link href="https://fonts.googleapis.com/css2?family=Rubik:wght@400;500;600;700;800&display=swap" rel="stylesheet">'
FILES = ["index.html", "style.css", "script.js"]
LINE_BUDGET = 1000


# --------------------------------------------------------------------------
# replaying the weeks
# --------------------------------------------------------------------------

def blocks_at(upto):
    """Ordered [(block_id, lines)] per file at the END of week `upto` (1-based)."""
    blocks = {name: [] for name in FILES}
    for week in course.WEEKS[:upto]:
        for kind, filename, block_id, lines in week["ops"]:
            existing = [i for i, (bid, _) in enumerate(blocks[filename]) if bid == block_id]
            if kind == "add":
                if existing:
                    raise SystemExit(f"week {week['n']}: block '{block_id}' already exists in {filename}; use SET")
                blocks[filename].append((block_id, lines))
            elif kind == "set":
                if not existing:
                    raise SystemExit(f"week {week['n']}: block '{block_id}' not in {filename} yet; use ADD")
                blocks[filename][existing[0]] = (block_id, lines)
            else:
                raise SystemExit(f"week {week['n']}: unknown op '{kind}'")
    # Sort into reading order, not authoring order. Week 2 teaches the submit
    # handler and week 3 the config it depends on, but the FILE has to put the
    # config first or students read a page that uses things it has not defined.
    out = {}
    for name in FILES:
        order = course.ORDER[name]
        unknown = [bid for bid, _ in blocks[name] if bid not in order]
        if unknown:
            raise SystemExit(f"{name}: block(s) {unknown} missing from course.ORDER[{name!r}]")
        out[name] = sorted(blocks[name], key=lambda pair: order.index(pair[0]))
    return out


def state_at(upto):
    """File contents as they stand at the END of week `upto`."""
    return {name: "\n".join(line for _, lines in ordered for line in lines)
            for name, ordered in blocks_at(upto).items()}


def block_spans(upto):
    """{filename: {block_id: (first_line, last_line)}} at the end of week `upto`.

    This is what lets the teacher guide say "type into script.js, lines 34-52"
    and be right. The numbers come out of the same replay that produces the
    milestone pages, so they cannot drift from what students actually see.
    """
    spans = {}
    for name, ordered in blocks_at(upto).items():
        spans[name] = {}
        line = 1
        for block_id, lines in ordered:
            spans[name][block_id] = (line, line + len(lines) - 1)
            line += len(lines)
    return spans


def week_ops(week):
    """{(filename, block_id): 'add' | 'set'} for one week."""
    return {(filename, block_id): kind for kind, filename, block_id, _ in week["ops"]}


MAX_STEP_LINES = 6   # the PixelPad workbooks' rule, and it holds up here too

# A trailing comma is NOT here on purpose: between two object properties it is a
# natural place to stop, and without that askAI's fetch call has no legal break
# point at all and lands as one fourteen-line wall.
_OPENERS = ("{", "(", "[", "=>", "&&", "||", "+")
_HEADERS = ("if", "for", "while", "function", "async", "try", "catch", "else", "return")


def _is_code(line):
    """Does this line count toward the six? Blanks and comments do not -
    they are the explanation, not the thing a student has to get right."""
    stripped = line.strip()
    return bool(stripped) and not stripped.startswith("//") and not stripped.startswith("/*")


def _depth_before(lines):
    """Nesting depth at the start of each line, relative to the block."""
    depth, out = 0, []
    for line in lines:
        code = line.split("//")[0]
        out.append(depth)
        depth += code.count("{") + code.count("(") + code.count("[")
        depth -= code.count("}") + code.count(")") + code.count("]")
    return out


def _can_break_before(lines, index, depths, css):
    """Never strand an `if` from its body, cut a half-finished expression, or
    start a piece on a bare closing brace."""
    if index == 0 or index >= len(lines):
        return False

    # A CSS rule is one idea to a beginner - never cut one open.
    if css and depths[index] > 0:
        return False

    # Starting a chunk with "}," or "})" reads as gibberish on a slide.
    following = lines[index].strip()
    if following.startswith(("}", ")", "]")):
        return False

    previous = lines[index - 1].strip()
    if not previous:
        return True                       # a blank line is the nicest place to stop
    if previous.endswith(_OPENERS):
        return False
    if previous.startswith("//"):
        return False                      # a comment belongs with what it describes
    first_word = previous.split("(")[0].split()[0] if previous.split() else ""
    if first_word in _HEADERS and not previous.endswith((";", "}")):
        return False
    return True


def chunk_block(lines, filename=""):
    """Split a block into pieces of at most six code lines each.

    Six is the cap the PixelPad workbooks use, and the reason is the same: past
    that a student is copying rather than following. Split points are chosen so
    an `if` never gets separated from the body it runs, a comment always travels
    with the code it explains, a CSS rule is never cut open, and no piece opens
    on a bare closing brace.

    An indivisible run with no legal break point may exceed six - the same
    allowance RULES.md makes for a single long `if` block.
    """
    css = filename.endswith(".css")
    depths = _depth_before(lines)
    chunks, current, code_count = [], [], 0
    for index, line in enumerate(lines):
        if code_count >= MAX_STEP_LINES and _can_break_before(lines, index, depths, css):
            # A chunk should not END on blank lines - they introduce whatever
            # comes next, so they travel forward rather than being dropped.
            carry = []
            while current and not current[-1].strip():
                carry.insert(0, current.pop())
            chunks.append(current)
            current, code_count = list(carry), 0
        current.append(line)
        if _is_code(line):
            code_count += 1
    if current:
        chunks.append(current)

    # Fold any all-blank chunk into its neighbour so no line is ever lost.
    merged = []
    for chunk in chunks:
        if merged and not any(line.strip() for line in chunk):
            merged[-1].extend(chunk)
        else:
            merged.append(list(chunk))
    assert [line for chunk in merged for line in chunk] == list(lines), \
        "chunking lost or reordered lines"
    return merged


def block_code(week_n, filename, block_id):
    """(first_line_number, lines, changed_line_numbers, removals) for one block."""
    start, end = block_spans(week_n)[filename][block_id]
    lines = state_at(week_n)[filename].split("\n")[start - 1:end]
    gone = removed_in_blocks(week_n)[filename].get(block_id, {})
    return start, lines, changed_lines(week_n)[filename], gone


def code_table(filename, start, lines, marks, ident=None, tag="", gone=None):
    """One styled snippet: file chip, line range, the code, changes accented.

    Deleted lines are drawn back in where they used to be, struck through, so a
    week that removes something can say so. Green alone cannot.
    """
    gone = gone or {}
    rows = []

    def removals(at):
        return "".join(
            f'<tr class="gone"><td class="ln">{at}</td>'
            f'<td class="src">{esc(text) or "&nbsp;"}</td></tr>'
            for text in gone.get(at, []))

    for offset, line in enumerate(lines):
        number = start + offset
        rows.append(removals(number))
        cls = ' class="new"' if number in marks else ""
        rows.append(f'<tr{cls}><td class="ln">{number}</td>'
                    f'<td class="src">{esc(line) or "&nbsp;"}</td></tr>')
    rows.append(removals(start + len(lines)))
    end = start + len(lines) - 1
    rng = f"line {start}" if start == end else f"lines {start}&ndash;{end}"
    copy = f'<button data-copy="{ident}">Copy</button>' if ident else ""
    table_id = f' id="{ident}"' if ident else ""
    return (f'<div class="snip"><div class="snip-head">'
            f'<span class="file">{esc(filename)}</span><span class="rng">{rng}</span>'
            f'<span class="tag">{tag}</span>{copy}</div>'
            f'<div class="code"><table{table_id}>{"".join(rows)}</table></div></div>')


def render_ask(ask):
    if not ask:
        return ""
    question, listening = ask
    return (f'<div class="say"><b>Ask the room:</b> {esc(question)}'
            f'<span class="listen">Listening for: {esc(listening)}</span></div>')


def placement(week_n, filename, block_id, kind):
    """Where in the file this block goes, said in a way you can act on.

    A teacher works down the guide, but the guide's order is the order things
    make sense in, not the order they sit in the file - week 5 teaches the
    error check inside askAI and only then the explain() function that lives
    above it. Line numbers alone do not help there, because they are the
    end-of-week numbers and the file is still growing. Naming the neighbour
    does: "goes just above function explain" is unambiguous at any point.
    """
    if kind == "set":
        start, end = block_spans(week_n)[filename][block_id]
        where = f"line {start}" if start == end else f"lines {start}&ndash;{end}"
        return (f"Already in <b>{esc(filename)}</b> at {where}. "
                f"<b>Only the green lines change</b> &mdash; everything else stays exactly as it is.")

    ordered = [bid for bid, _ in blocks_at(week_n)[filename]]
    index = ordered.index(block_id)
    if index == 0:
        return f"New. Goes at the very top of <b>{esc(filename)}</b>, above everything else."
    if index == len(ordered) - 1:
        return f"New. Goes at the end of <b>{esc(filename)}</b>."
    following = dict(blocks_at(week_n)[filename])[ordered[index + 1]]
    anchor = next((line.strip() for line in following if line.strip()), "")
    # Trim on a word boundary so the anchor never ends mid-word.
    trimmed = esc(anchor) if len(anchor) <= 62 else esc(anchor[:62].rsplit(" ", 1)[0]) + "&hellip;"
    return (f"New. Goes in <b>{esc(filename)}</b>, just above the line "
            f"<code>{trimmed}</code>")


def render_step(week, beat):
    """A typing beat, broken into pieces of at most six code lines.

    Each piece gets its own sentence of explanation, so the room stops and
    talks roughly every half-dozen lines instead of copying twenty in silence.
    """
    filename, block_id = beat["file"], beat["block"]
    start, lines, marks, gone = block_code(week["n"], filename, block_id)
    chunks = chunk_block(lines, filename)
    notes = beat.get("notes") or []

    if len(notes) != len(chunks):
        offsets, cursor = [], start
        for chunk in chunks:
            offsets.append(f"{cursor}-{cursor + len(chunk) - 1}")
            cursor += len(chunk)
        raise SystemExit(
            f"week {week['n']} STEP {filename}:{block_id} splits into {len(chunks)} "
            f"chunk(s) at lines {', '.join(offsets)}, but you wrote {len(notes)} note(s). "
            f"Add one short note per chunk."
        )

    kind = week_ops(week).get((filename, block_id))
    tag = {"add": "new this week", "set": "replaces what is there"}.get(kind, "already written")
    pieces, cursor = [], start
    for index, (chunk, note) in enumerate(zip(chunks, notes)):
        last = cursor + len(chunk) - 1
        final = index == len(chunks) - 1
        touched = any(number in marks for number in range(cursor, last + 1))
        # A line removed from the end of a block sits at the number of whatever
        # follows it, so only the last chunk reaches past its own end - or a
        # deletion at a chunk boundary is drawn twice.
        here = {at: text for at, text in gone.items()
                if cursor <= at <= (last + 1 if final else last)}
        if kind == "set" and not touched and not here:
            # Do not print twenty lines a student already has just to reach the
            # six that changed. Say what is there, say to leave it, move on.
            span = f"Line {cursor}" if cursor == last else f"Lines {cursor}&ndash;{last}"
            pieces.append(f'<p class="unchanged"><b>{span} do not change</b> &mdash; '
                          f'skip past them. <span class="muted">{note}</span></p>')
        else:
            # `here` is why the skip above also tests for removals. Week 4's
            # only edit to the models block is that a line LEAVES it, and
            # nothing in the block goes green - so on the count of changed
            # lines alone this printed "lines 29-38 do not change" for the one
            # block the week is about.
            pieces.append(code_table(filename, cursor, chunk, marks, tag=tag, gone=here))
            pieces.append(f'<p class="chunk-note">{note}</p>')
        cursor += len(chunk)

    at = f'<span class="at">{esc(beat["at"])}</span>' if beat.get("at") else ""
    lead = f'<p class="step-lead">{placement(week["n"], filename, block_id, kind)}</p>'
    return (f'<section class="beat step"><h4>{at}{esc(beat["title"])}</h4>{lead}'
            f'{"".join(pieces)}{render_ask(beat.get("ask"))}</section>')


def render_flow(week):
    out = []
    for beat in week["flow"]:
        if beat["kind"] == "step":
            out.append(render_step(week, beat))
        else:
            at = f'<span class="at">{esc(beat["at"])}</span>' if beat.get("at") else ""
            body = "".join(f"<p>{para}</p>" for para in beat["body"])
            out.append(f'<section class="beat"><h4>{at}{esc(beat["title"])}</h4>'
                       f'{body}{render_ask(beat.get("ask"))}</section>')
    return "".join(out)


def changed_lines(upto):
    """Line numbers (1-based, per file) that are new or changed this week.

    Computed by diffing the replayed state rather than tracked by hand, so a
    week that rewrites an earlier function highlights correctly without the
    author having to say so.
    """
    before = state_at(upto - 1) if upto > 1 else {name: "" for name in FILES}
    after = state_at(upto)
    marks = {}
    for name in FILES:
        old = before[name].split("\n") if before[name] else []
        new = after[name].split("\n") if after[name] else []
        hits = set()
        matcher = difflib.SequenceMatcher(None, old, new, autojunk=False)
        for tag, _i1, _i2, j1, j2 in matcher.get_opcodes():
            if tag in ("insert", "replace"):
                hits.update(range(j1 + 1, j2 + 1))
        marks[name] = hits
    return marks


def line_sig(text):
    """A line's CODE signature: whitespace collapsed and any comment stripped, so
    two lines match when only their indentation or comments differ. It lets us
    tell a line that truly went away from one that merely moved or lost a comment.
    A comment here is set off by two-plus spaces or starts the line, so a URL's
    '://' survives. Deliberately simple, matching the light comment style the
    decks use. A whole-line comment returns '' - it has no code to match on."""
    text = re.sub(r"\s{2,}//.*$", "", text)   # a trailing, space-aligned comment
    text = re.sub(r"^\s*//.*$", "", text)     # a whole-line comment
    return " ".join(text.split())


def removed_in_blocks(upto):
    """{filename: {block_id: {line_in_new_file: [text, ...]}}} - deletions,
    attributed to the block they were removed FROM.

    Green-only highlighting cannot show a deletion: a removed line is not in
    the new file to colour. Week 4 drops the listModels() call and the page had
    no way to say so, leaving a student comparing their file against the page
    with one extra line and no explanation.

    Diffed per BLOCK, not per file. Across a whole file, difflib pairs a
    deletion with whatever was inserted nearby and calls the result "replace" -
    week 4's dropped call came back as replace(old[39:40] -> new[39:58]),
    tangled up with the askAI block appearing above it. Inside a single block
    there is nothing to tangle with, so a delete is unambiguously a delete.

    The block id is kept rather than only the line number, because a line
    dropped from the END of a block sits at the first line number of whatever
    follows it. Week 4's removal is at position 39, which is one past `models`
    AND the first line of `ask` - a range test cannot tell those apart, and
    showing it on both is wrong twice over.
    """
    if upto <= 1:
        return {name: {} for name in FILES}
    before = {name: dict(blocks) for name, blocks in blocks_at(upto - 1).items()}
    after = blocks_at(upto)
    spans = block_spans(upto)
    gone = {}
    for name in FILES:
        per_block = {}
        for block_id, new_lines in after[name]:
            old_lines = before[name].get(block_id)
            if old_lines is None or old_lines == new_lines:
                continue
            block_start = spans[name][block_id][0]
            new_sigs = {sig for sig in (line_sig(l) for l in new_lines) if sig}
            matcher = difflib.SequenceMatcher(None, old_lines, new_lines, autojunk=False)
            for tag, i1, i2, j1, j2 in matcher.get_opcodes():
                # A line the student must take out - whether it simply vanished
                # (a delete) or was swapped for something else (a replace). The
                # trap is a block that is only re-indented or re-commented:
                # difflib calls every such line "replace", and striking the ones
                # that really survive - week 5 wraps three lines in try/catch,
                # unchanged but indented - would bury the line that truly went
                # AND make them retype code that did not change. So compare by
                # CODE SIGNATURE (whitespace collapsed, comments dropped) and
                # strike a line only when nothing with its signature remains.
                if tag not in ("delete", "replace"):
                    continue
                dropped = [line for line in old_lines[i1:i2]
                           if line.strip() and line_sig(line) not in new_sigs]
                if dropped:
                    per_block.setdefault(block_id, {}).setdefault(
                        block_start + j1, []).extend(dropped)
        gone[name] = per_block
    return gone


def removed_lines(upto):
    """{filename: {line_in_new_file: [text, ...]}} - the same deletions, keyed
    only by position, for the whole-file views where no block is in play."""
    flat = {}
    for name, per_block in removed_in_blocks(upto).items():
        at = {}
        for positions in per_block.values():
            for position, texts in positions.items():
                at.setdefault(position, []).extend(texts)
        flat[name] = at
    return flat


def line_count(files):
    return sum(len(text.split("\n")) for text in files.values() if text)


# --------------------------------------------------------------------------
# shared page furniture
# --------------------------------------------------------------------------

def esc(text):
    """HTML-escape, and force the result to pure ASCII.

    Every page this builder writes is 7-bit clean. Non-ASCII becomes a numeric
    character reference, which renders identically but survives being opened by
    something that guesses the encoding wrong - Word, an old editor, a
    PowerShell round-trip. These files get printed, emailed and dropped into
    Google Drive, so that is a real risk, and a mojibake page in front of a
    class is worse than a slightly longer file.
    """
    return html.escape(str(text)).encode("ascii", "xmlcharrefreplace").decode("ascii")


def guard(tool="ai101", up="../"):
    """The site access gate. `up` is the path back to the site root - the slide
    decks sit one folder deeper than everything else this builder writes."""
    return (
        '<script>window.UTG_GUARD="%(up)s";window.UTG_TOOL="%(tool)s";'
        "document.write('<script src=\"%(up)sclass-codes.js?t='+Date.now()+'\"><\\/script>');</script>\n"
        '<script src="%(up)sguard.js"></script>' % {"tool": tool, "up": up}
    )


CSS = """
:root{--brand:#01aefd;--brand-dark:#0294d8;--brand-ink:#0a6299;--brand-tint:#e7f7ff;
--gold:#ffd633;--ink:#1f2a37;--muted:#6b7787;--bg:#eef2f7;--surface:#fff;--border:#dbe3ec;
--add:#e8f8ec;--add-edge:#37a95d;}
*{box-sizing:border-box}
body{margin:0;background:var(--bg);color:var(--ink);font:16px/1.6 Rubik,system-ui,-apple-system,"Segoe UI",sans-serif}
a{color:var(--brand);text-decoration:none}
.wrap{max-width:1080px;margin:0 auto;padding:26px 22px 80px}
header.site{display:flex;align-items:center;gap:12px;padding:16px 22px;background:var(--surface);border-bottom:1px solid var(--border)}
header.site img{height:26px}
header.site .slash{color:#b8c8ce}
header.site strong{font-size:15px;letter-spacing:.02em}
.eyebrow{color:var(--brand);font-size:12px;font-weight:800;letter-spacing:.12em;text-transform:uppercase;margin:0 0 10px}
h1{font-size:34px;margin:0 0 10px;line-height:1.2}
h2{font-size:23px;margin:34px 0 12px}
h3{font-size:17px;margin:22px 0 8px}
.lead{color:var(--muted);font-size:17px;max-width:70ch}
.card{background:var(--surface);border:1px solid var(--border);border-radius:9px;padding:20px;margin:16px 0}
.grid{display:grid;gap:12px;grid-template-columns:repeat(auto-fill,minmax(230px,1fr))}
.wk-card{display:block;background:var(--surface);border:1px solid var(--border);border-radius:9px;padding:16px;color:inherit}
.wk-card:hover{border-color:var(--brand)}
.wk-card .n{color:var(--brand);font-size:12px;font-weight:800;letter-spacing:.1em;text-transform:uppercase}
.wk-card h3{margin:6px 0 4px;font-size:16px}
.wk-card p{margin:0;color:var(--muted);font-size:13px;line-height:1.5}
.pill{display:inline-block;background:var(--brand-tint);color:var(--brand-ink);border-radius:20px;padding:5px 12px;font-size:12px;font-weight:700;margin:0 6px 6px 0}
.tabs{display:flex;gap:0;border-bottom:1px solid var(--border);margin-top:6px}
.tab{border:0;border-right:1px solid var(--border);background:#f7fafb;color:#52707a;padding:10px 14px;font:700 12px Rubik,sans-serif;cursor:pointer}
.tab.active{background:var(--surface);color:var(--brand);box-shadow:inset 0 -2px var(--brand)}
/* ---- code ---- */
.snip{margin:14px 0;border:1px solid #223f49;border-radius:9px;overflow:hidden;background:#0f1b21}
.snip-head{display:flex;align-items:center;gap:10px;padding:9px 14px;background:#16303a;
  color:#9fc4cf;font-size:12px;letter-spacing:.02em}
.snip-head .file{color:#8fd7f0;font:700 12px ui-monospace,SFMono-Regular,Menlo,Consolas,monospace}
.snip-head .rng{color:#6d8c97}
.snip-head .tag{margin-left:auto;color:#6d8c97}
.snip-head button{border:1px solid #33586a;background:transparent;color:#a8ccd8;border-radius:5px;
  padding:4px 10px;font:600 11px Rubik,sans-serif;cursor:pointer}
.snip-head button:hover{border-color:#5f93a8;color:#dff0f6}
.snip .code{overflow:auto;margin:0;padding:10px 0}
.snip table{border-collapse:collapse;width:100%;
  font:13px/1.85 ui-monospace,SFMono-Regular,Menlo,Consolas,monospace}
.snip td{padding:0;white-space:pre;vertical-align:top}
.snip td.ln{width:1%;padding:0 14px 0 16px;text-align:right;color:#44646f;user-select:none}
.snip td.src{color:#d7e6ea;padding-right:18px}
/* new lines get a quiet tint and an accent bar, not a solid block of colour */
.snip tr.new td.src{background:rgba(63,186,110,.12)}
.snip tr.new td.ln{background:rgba(63,186,110,.12);color:#79d3a0;box-shadow:inset 3px 0 #3fba6e}
/* a line this week removes - shown where it used to be, struck through */
.snip tr.gone td.src{background:rgba(201,74,58,.10);color:#e39289;text-decoration:line-through;
  text-decoration-color:rgba(227,146,137,.65)}
.snip tr.gone td.ln{background:rgba(201,74,58,.10);color:#e39289;box-shadow:inset 3px 0 #c94a3a}

.teach-row{margin:14px 0 18px}
.pill.ghost{background:var(--surface);color:var(--brand-ink);border:1px solid var(--border)}
.pill.ghost:hover{border-color:var(--brand);background:var(--brand-tint)}
@media print{.teach-row{display:none}}

/* ---- week anchors and the jump bar ---- */
.chapter{scroll-margin-top:74px}
.jump{display:flex;flex-wrap:wrap;gap:6px;align-items:center;margin:0 0 18px}
.jump span{font-size:12px;font-weight:800;letter-spacing:.08em;text-transform:uppercase;
color:var(--muted);margin-right:4px}
.jump a{display:grid;place-items:center;min-width:30px;height:30px;padding:0 7px;
background:var(--surface);border:1px solid var(--border);border-radius:7px;
font-size:13px;font-weight:700;color:var(--brand-ink);text-decoration:none}
.jump a:hover{border-color:var(--brand);background:var(--brand-tint)}
/* Opened inside the classroom the page is already inside a titled panel, so
   the site header, the intro and the print button would only be noise. */
body.embedded header.site,body.embedded .printbtn,body.embedded .jump,
body.embedded footer{display:none}
body.embedded .wrap{padding-top:8px}
@media print{.jump{display:none}}

/* ---- lesson flow ---- */
.beat{margin:0 0 4px;padding:14px 0 14px 20px;border-left:2px solid var(--border)}
.beat:hover{border-left-color:#c3d3dd}
.beat.step{border-left-color:var(--brand)}
.beat h4{margin:0 0 8px;font-size:16.5px;display:flex;align-items:baseline;gap:11px}
.beat h4 .at{color:var(--brand);font:800 11.5px/1 Rubik,sans-serif;letter-spacing:.09em;
  min-width:40px;flex:none;padding-top:2px}
.beat p{margin:0 0 9px;line-height:1.68;max-width:74ch}
.beat p:last-child{margin-bottom:0}
.chunk-note{margin:2px 0 0;color:#3f5764;font-size:14.5px;line-height:1.62;max-width:74ch}
.unchanged{margin:10px 0;padding:9px 13px;border-left:3px solid #c3d3dd;background:#f4f7f9;
  border-radius:0 6px 6px 0;color:#4a5f6b;font-size:13.5px;line-height:1.55;max-width:74ch}
.step-lead{color:var(--muted);font-size:13.5px;margin:0 0 4px}
.note{border-left:3px solid var(--brand);background:var(--brand-tint);padding:12px 14px;margin:14px 0;font-size:14px;border-radius:0 6px 6px 0}
.warn{border-left:3px solid #e9952a;background:#fff8e9;padding:12px 14px;margin:14px 0;font-size:14px;border-radius:0 6px 6px 0}
.bonus{border-left:3px solid var(--gold);background:#fffbea;padding:12px 14px;margin:14px 0;font-size:14px;border-radius:0 6px 6px 0}
.muted{color:var(--muted)}
.nav{display:flex;justify-content:space-between;gap:12px;margin-top:34px;padding-top:18px;border-top:1px solid var(--border)}
footer{color:var(--muted);font-size:13px;padding:24px 22px;text-align:center}
table.plan{width:100%;border-collapse:collapse;margin:10px 0;font-size:14px}
table.plan th,table.plan td{border:1px solid var(--border);padding:8px 10px;text-align:left;vertical-align:top}
table.plan th{background:var(--brand-tint);color:var(--brand-ink);font-size:12px;text-transform:uppercase;letter-spacing:.06em}
table.plan td.t{width:70px;white-space:nowrap;color:var(--muted);font-weight:700}
ul.tight li{margin-bottom:5px}
.say{background:#f6f2fd;border:1px solid #e2d8f6;padding:11px 14px;margin:11px 0 0;
  border-radius:7px;font-size:14.2px;line-height:1.6;max-width:74ch}
.say b{color:#5b3fa0}
.say .listen{display:block;margin-top:5px;color:#6b6280;font-size:13.2px}
@media print{
  @page{size:letter;margin:14mm}
  body{background:#fff}
  header.site,.nav,.snip-head button,.noprint{display:none}
  .beat{break-inside:avoid;page-break-inside:avoid}
  .snip{break-inside:avoid;page-break-inside:avoid}
  .chapter{page-break-before:always;page-break-inside:avoid}
  .chapter:first-of-type{page-break-before:avoid}
  .card,.note,.warn,.bonus{page-break-inside:avoid}
}
html:not(.utg-can-print) .printbtn{display:none}
@media print{html:not(.utg-can-print) body{display:none}}
"""

COPY_JS = """
document.querySelectorAll('[data-copy]').forEach(function(btn){
  btn.addEventListener('click', function(){
    var pre = document.getElementById(btn.getAttribute('data-copy'));
    var text = Array.prototype.map.call(pre.querySelectorAll('td.src'), function(td){ return td.textContent; }).join('\\n');
    navigator.clipboard.writeText(text).then(function(){
      var was = btn.textContent; btn.textContent = 'Copied'; setTimeout(function(){ btn.textContent = was; }, 1400);
    });
  });
});
document.querySelectorAll('[data-tabs]').forEach(function(group){
  group.querySelectorAll('.tab').forEach(function(tab){
    tab.addEventListener('click', function(){
      group.querySelectorAll('.tab').forEach(function(t){ t.classList.remove('active'); });
      tab.classList.add('active');
      var owner = group.parentElement;
      owner.querySelectorAll('[data-pane]').forEach(function(pane){
        pane.style.display = pane.getAttribute('data-pane') === tab.getAttribute('data-for') ? '' : 'none';
      });
    });
  });
});
"""


def page(title, body, extra_js="", tool="ai101"):
    return f"""<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{esc(title)} &middot; UTG Academy</title>
{FONT}
<style>{CSS}</style>
{guard(tool)}
<header class="site"><a href="../"><img src="{LOGO}" alt="UTG Academy"></a><span class="slash">/</span><strong>AI101</strong></header>
{body}
<footer>&copy; 2026 UTG Academy</footer>
<script>{COPY_JS}{extra_js}</script>
"""


def code_block(filename, text, marks, ident, gone=None):
    """A whole file rendered with line numbers, this week's changes accented."""
    label = "green = new" + (", struck through = deleted" if gone else "") + " this week"
    return code_table(filename, 1, text.splitlines(), marks,
                      ident=ident, tag=label, gone=gone)


def files_view(upto, ident_prefix):
    """The three files as clickable tabs, with the current week highlighted."""
    files = state_at(upto)
    marks = changed_lines(upto)
    gone = removed_lines(upto)
    tabs = "".join(
        f'<button class="tab{" active" if i == 0 else ""}" data-for="{ident_prefix}-{i}">{esc(name)}</button>'
        for i, name in enumerate(FILES)
    )
    panes = "".join(
        f'<div data-pane="{ident_prefix}-{i}"{"" if i == 0 else ' style="display:none"'}>'
        f'{code_block(name, files[name], marks[name], f"{ident_prefix}-code-{i}", gone[name])}</div>'
        for i, name in enumerate(FILES)
    )
    return f'<div class="tabs" data-tabs>{tabs}</div>{panes}'


# --------------------------------------------------------------------------
# the four deliverables
# --------------------------------------------------------------------------

def build_index():
    cards = "".join(
        f'<a class="wk-card" href="week-{w["n"]:02d}.html"><span class="n">Week {w["n"]}</span>'
        f'<h3>{esc(w["title"])}</h3><p>{esc(w["big_idea"])}</p></a>'
        for w in course.WEEKS
    )
    body = f"""<div class="wrap">
<p class="eyebrow">15-week course</p>
<h1>{esc(course.COURSE_TITLE)}</h1>
<p class="lead">{esc(course.COURSE_BLURB)}</p>
<div class="card">
  <h3 style="margin-top:0">What you are building</h3>
  <p class="muted">{esc(course.PROJECT_BLURB)}</p>
  <p><a class="pill" href="teacher.html">Teacher curriculum</a>
     <a class="pill" href="workbook.html">Student homework book</a>
     <a class="pill" href="../classroom/">Open the code editor</a></p>
  <p class="muted" style="margin:10px 0 0;font-size:13px">Each week page carries that week&rsquo;s
  slides and links straight to its lesson plan. In the classroom, opening a week gives the teacher
  both without leaving the room.</p>
</div>
<div class="warn"><h3 style="margin-top:0">Before you start &mdash; read this</h3>{course.DISCLAIMER}</div>
<h2>The weeks</h2>
<div class="grid">{cards}</div>
</div>"""
    write("index.html", page(course.COURSE_TITLE, body))


def build_weeks():
    for index, week in enumerate(course.WEEKS):
        n = week["n"]
        prev = f'<a href="week-{n-1:02d}.html">&larr; Week {n-1}</a>' if n > 1 else "<span></span>"
        nxt = f'<a href="week-{n+1:02d}.html">Week {n+1} &rarr;</a>' if n < len(course.WEEKS) else "<span></span>"
        concepts = "".join(f'<span class="pill">{esc(c)}</span>' for c in week["new_concepts"])
        total = line_count(state_at(n))
        dropped = sum(len(v) for name in FILES for v in removed_lines(n)[name].values())
        legend = (f"The green lines are what is new since week {n-1}."
                  if n > 1 else "Everything here is new - this is week 1.")
        if dropped:
            legend += (" The struck-through line is one you DELETE this week - take it out."
                       if dropped == 1 else
                       f" The {dropped} struck-through lines are ones you DELETE this week"
                       " - take them out.")
        if n > 1:
            legend += " Everything else you already had."
        bonus = ""
        if week.get("bonus"):
            bonus = (
                f'<div class="bonus"><strong>Finished early? {esc(week["bonus"]["title"])}</strong>'
                f'<p style="margin:6px 0 0">{esc(week["bonus"]["body"])}</p>'
                '<p class="muted" style="margin:6px 0 0;font-size:13px">Nothing in a later week depends on this, '
                'so it is safe to skip and safe to keep.</p></div>'
            )
        body = f"""<div class="wrap">
<p class="eyebrow">Week {n} of {len(course.WEEKS)}</p>
<h1>{esc(week["title"])}</h1>
<p class="lead">{esc(week["big_idea"])}</p>
<p>{concepts}</p>
<p class="teach-row"><a class="pill" href="slides/week-{n:02d}.html">Slides for week {n}</a>
   <a class="pill ghost" href="slides/week-{n:02d}.pptx">.pptx</a>
   <a class="pill ghost" href="teacher.html#week-{n}">Lesson plan for week {n}</a></p>
<div class="note"><strong>Where you should be by the end of this week.</strong>
{legend}
Your project is now {total} lines.</div>
{bonus}
<h2>Your code after week {n}</h2>
{files_view(n, f"w{n}")}
<div class="nav">{prev}{nxt}</div>
</div>"""
        write(f"week-{n:02d}.html", page(f"Week {n} · {week['title']}", body))


def build_teacher():
    sections = []
    for week in course.WEEKS:
        spans = block_spans(week["n"])
        # An at-a-glance list of every edit the week makes, straight from the ops.
        typed = "".join(
            '<tr><td><b>{f}</b></td><td>{r}</td><td>{k}</td></tr>'.format(
                f=esc(filename),
                r=("line {0}".format(spans[filename][block_id][0])
                   if spans[filename][block_id][0] == spans[filename][block_id][1]
                   else "lines {0}&ndash;{1}".format(*spans[filename][block_id])),
                k="type it in fresh" if kind == "add" else "replace what is already there")
            for kind, filename, block_id, _ in week["ops"]
        )
        typed_table = (
            '<div class="card"><strong>Everything typed this week</strong>'
            '<table class="plan"><tr><th>File</th><th>Where</th><th>What to do</th></tr>'
            f'{typed}</table>'
            '<p class="muted" style="margin:8px 0 0;font-size:13px">Line numbers are as the file '
            'stands at the END of this week, which is what the student sees on the week page. '
            'Work top to bottom and they will line up.</p></div>'
            if typed else '<div class="card"><strong>No new code this week.</strong></div>'
        )
        errors = "".join(
            f"<li><strong>{esc(sym)}</strong> &mdash; {esc(fix)}</li>" for sym, fix in week["errors"]
        )
        bonus = ""
        if week.get("bonus"):
            bonus = (f'<div class="bonus"><strong>Bonus for fast finishers: {esc(week["bonus"]["title"])}</strong>'
                     f'<p style="margin:6px 0 0">{esc(week["bonus"]["body"])}</p></div>')
        sections.append(f"""<section class="chapter" id="week-{week["n"]}">
<h2>Week {week["n"]} &middot; {esc(week["title"])}</h2>
<p class="lead">{esc(week["big_idea"])}</p>
<div class="card"><strong>They leave today able to:</strong>
<ul class="tight">{"".join(f"<li>{esc(o)}</li>" for o in week["objectives"])}</ul></div>
{typed_table}
<h3>How the hour goes</h3>
{render_flow(week)}
<h3>What will go wrong</h3>
<ul class="tight">{errors}</ul>
{bonus}
<h3>Homework set today</h3>
<ul class="tight">{"".join(f"<li>{esc(c['task'])}</li>" for c in week["homework"])}</ul>
</section>""")
    jump = "".join(
        f'<a href="#week-{week["n"]}">{week["n"]}</a>' for week in course.WEEKS
    )
    body = f"""<div class="wrap">
<p class="eyebrow">Teacher curriculum</p>
<h1>{esc(course.COURSE_TITLE)}</h1>
<p class="lead">Fifteen one-hour lessons. Every hour is built so you are talking for well under half of it.</p>
<nav class="jump" aria-label="Jump to a week"><span>Week</span>{jump}</nav>
<button class="printbtn pill" onclick="window.print()">Print to PDF</button>
<div class="warn"><h3 style="margin-top:0">Say this in week 1</h3>{course.DISCLAIMER}</div>
<div class="card">
  <h3 style="margin-top:0">How to run this course</h3>
  {course.TEACHER_PREAMBLE}
</div>
{"".join(sections)}
</div>"""
    # ?embed=1 strips the page furniture. The browser handles #week-N on its
    # own, but only if the element exists when it looks - so nudge it on load
    # too, and again on a hash change, which is how the panel switches weeks
    # without reloading the frame.
    embed_js = (
        "if(new URLSearchParams(location.search).has('embed'))"
        "document.body.classList.add('embedded');"
        "function utgJump(){var id=location.hash.slice(1);if(!id)return;"
        "var el=document.getElementById(id);if(el)el.scrollIntoView();}"
        "addEventListener('load',utgJump);addEventListener('hashchange',utgJump);"
    )
    write("teacher.html", page("Teacher curriculum", body, extra_js=embed_js))


def build_workbook():
    chapters = []
    for week in course.WEEKS:
        challenges = "".join(
            f'<div class="card"><strong>Challenge {i}. {esc(c["task"])}</strong>'
            f'<p style="margin:8px 0 0">{esc(c["detail"])}</p>'
            f'<p class="muted" style="margin:8px 0 0;font-size:13px"><strong>Done when:</strong> {esc(c["done"])}</p></div>'
            for i, c in enumerate(week["homework"], start=1)
        )
        recap = "".join(f"<li>{esc(r)}</li>" for r in week["recap"])
        chapters.append(f"""<section class="chapter">
<p class="eyebrow">Chapter {week["n"]}</p>
<h2>{esc(week["title"])}</h2>
<p class="lead">{esc(week["big_idea"])}</p>
<h3>What you learned</h3>
<ul class="tight">{recap}</ul>
<h3>Words to know</h3>
<p>{"".join(f'<span class="pill">{esc(c)}</span>' for c in week["new_concepts"])}</p>
<h3>Challenges</h3>
{challenges}
<h3>Stuck?</h3>
<p class="muted">Open <a href="week-{week["n"]:02d}.html">week {week["n"]}'s code</a> and compare it with yours line by line.
The green lines are the ones added that week.</p>
</section>""")
    body = f"""<div class="wrap">
<p class="eyebrow">Student homework book</p>
<h1>{esc(course.COURSE_TITLE)}</h1>
<p class="lead">One chapter per week. Do the challenges on your own project - there is no separate file to make.</p>
<button class="printbtn pill" onclick="window.print()">Print to PDF</button>
<div class="warn"><h3 style="margin-top:0">Read this first</h3>{course.DISCLAIMER}</div>
{"".join(chapters)}
</div>"""
    write("workbook.html", page("Student homework book", body))


CODE_LINES_PER_SLIDE = 15   # beyond this the type is too small to read from the back


def code_chunks(week_n, refs):
    """[(filename, first_line, lines, changed)] split into slide-sized pieces.

    A thirty-line function will not fit on one readable slide, so it becomes two
    slides rather than shrinking to eight point. Line numbers keep running, so
    a student can always match the slide against their own file.
    """
    chunks = []
    for filename, block_id in refs:
        start, lines, marks, gone = block_code(week_n, filename, block_id)
        pieces = list(range(0, len(lines), CODE_LINES_PER_SLIDE)) or [0]
        for offset in pieces:
            piece = lines[offset:offset + CODE_LINES_PER_SLIDE]
            first = start + offset
            last = first + len(piece) - 1
            # A deletion belongs to the piece holding the line it used to sit
            # at. The final piece also takes anything past its end, so a block
            # that ends by removing its last line still says so.
            final = offset == pieces[-1]
            here = {at: text for at, text in gone.items()
                    if first <= at <= (last + 1 if final else last)}
            chunks.append((filename, first, piece, marks, here))
    return chunks


def notes_for(week_n, filename, block):
    """Per-line notes for a block, this week. A block that changes across weeks
    (ask, submit) keeps a (week, file, block) entry so each week's arrangement
    gets its own notes; blocks that appear once use the plain (file, block) key."""
    return (course.LINE_NOTES.get((week_n, filename, block))
            or course.LINE_NOTES.get((filename, block)))


# A tiny bridge injected at the TOP of a checkpoint page, so it runs before the
# project's own script. It forwards console output, errors and fetch outcomes up
# to the deck, which shows them in a console under the output. That is how a
# silent failure becomes visible: week 4 has no error handling, so a rejected
# request just leaves "thinking..." on screen - here the console says why (a 401
# bad key, or a blocked/failed fetch off the school network).
CP_BRIDGE = """<script>(function () {
  if (window.parent === window) return;            // only when embedded in the deck
  function say(level, parts) {
    try {
      var text = Array.prototype.map.call(parts, function (a) {
        if (a instanceof Error) return a.message;
        if (a && typeof a === 'object') { try { return JSON.stringify(a); } catch (e) { return String(a); } }
        return String(a);
      }).join(' ');
      window.parent.postMessage({ utg: 'cp-log', level: level, text: text }, '*');
    } catch (e) {}
  }
  ['log', 'info', 'warn', 'error'].forEach(function (name) {
    var orig = console[name];
    console[name] = function () { say(name, arguments); if (orig) orig.apply(console, arguments); };
  });
  window.addEventListener('error', function (e) {
    say('error', [e.message + (e.filename ? ' (' + e.filename + ':' + e.lineno + ')' : '')]);
  });
  window.addEventListener('unhandledrejection', function (e) {
    var r = e.reason; say('error', ['Uncaught (in promise) ' + (r && r.message ? r.message : r)]);
  });
  var f = window.fetch;
  if (f) window.fetch = function () {
    var url = (arguments[0] && arguments[0].url) || arguments[0];
    say('net', ['\\u2192 ' + url]);
    return f.apply(this, arguments).then(function (res) {
      say(res.ok ? 'net' : 'error', [(res.ok ? '\\u2190 ' : '\\u2717 ') + res.status + ' ' + url]);
      return res;
    }, function (err) {
      say('error', ['\\u2717 could not reach ' + url + ' \\u2014 ' + (err && err.message || err)]);
      throw err;
    });
  };
})();</script>
"""


def checkpoint_page(week_n):
    """The runnable project at the end of week_n: index.html with its CSS and JS
    inlined, so a deck slide can render exactly what a student should be seeing,
    with a console bridge on top so the deck can show what it logs and where it
    fails. Weeks 1-2 run fully in the slide (no network); weeks that call the AI
    show the interface, and their reply needs the classroom key and network."""
    files = state_at(week_n)
    html = files.get("index.html", "")
    html = html.replace('<link rel="stylesheet" href="style.css">',
                        "<style>\n" + files.get("style.css", "") + "\n</style>")
    html = html.replace('<script src="script.js"></script>',
                        "<script>\n" + files.get("script.js", "") + "\n</script>")
    return CP_BRIDGE + html


def balance_quiz(quiz):
    """Move the correct option off "always A". Quizzes are authored with the
    right answer first (easy to read); this deterministically reseats it at a
    position derived from the question text, so across a deck the answers land
    on A/B/C/D roughly evenly and the same question always renders the same way
    (stable across the .html, the .pptx, and rebuilds). The other options keep
    their relative order."""
    opts = quiz["options"]
    n = len(opts)
    correct = opts[quiz["answer"]]
    target = int(hashlib.md5(quiz["q"].encode("utf-8")).hexdigest(), 16) % n
    others = [o for i, o in enumerate(opts) if i != quiz["answer"]]
    new = []
    oi = 0
    for pos in range(n):
        if pos == target:
            new.append(correct)
        else:
            new.append(others[oi])
            oi += 1
    return {**quiz, "options": new, "answer": target}


def slide_plan(week, seen):
    """Ordered slide descriptors for a week's body (the caller adds the title
    slide). Mutates `seen` - the concept keys already introduced course-wide -
    so a concept is explained only the first time it appears.

    An EXPANDED week gets a short explainer slide before the first line that uses
    a new HTML tag / CSS property / JavaScript idea, and one code slide per line
    with a note. Other weeks keep the chunked format, unchanged."""
    expanded = week["n"] in getattr(course, "EXPANDED_WEEKS", set())
    out = []
    for spec in week["slides"]:
        if spec.get("checkpoint"):
            # A "run it and check" slide: the deck can render the project as it
            # stands at this point, so the class can compare their screen to it.
            out.append(("checkpoint", {"week": spec.get("checkpoint_week", week["n"]),
                                       "title": spec.get("title", "Checkpoint"),
                                       "say": spec.get("say", "")}))
            continue
        out.append(("concept", {"eyebrow": "", "title": spec["title"],
                                 "sub": spec.get("sub", ""), "bullets": spec.get("bullets", [])}))
        refs = spec.get("code")
        if not refs:
            continue
        if not expanded:
            chunks = code_chunks(week["n"], refs)
            for index, (filename, start, lines, marks, gone) in enumerate(chunks, start=1):
                out.append(("chunk", {"file": filename, "start": start, "lines": lines,
                                      "marks": marks, "gone": gone, "part": index, "parts": len(chunks)}))
            continue
        eyebrow_for = {"html": "HTML", "css": "CSS", "js": "JavaScript"}
        delete_notes = getattr(course, "DELETE_NOTES", {})
        for filename, block in refs:
            start, lines, marks, gone = block_code(week["n"], filename, block)
            notes = notes_for(week["n"], filename, block)
            shown = 0

            def emit_deletes(at, context):
                nonlocal shown
                # A new line landing at this spot means the old one is being
                # swapped out, not just cut - say so, and the green "type this"
                # slide for its replacement comes right after.
                replacing = at in marks
                default = ("This line changes - take the old one out; the new "
                           "version is next." if replacing
                           else "Delete this line - take it out.")
                for dropped in gone.get(at, []):
                    out.append(("delete", {"file": filename, "start": start,
                                           "lines": context, "dropped": dropped, "at": at,
                                           "replacing": replacing,
                                           "note": delete_notes.get((filename, block), default)}))
                    shown += 1

            for i, line in enumerate(lines):
                lineno = start + i
                emit_deletes(lineno, lines[:i])   # a line removed from this spot
                if not line.strip():
                    continue
                # Only a line that is NEW or CHANGED this week gets its own "type
                # this" slide; a line carried over from an earlier week just shows
                # as dimmed context. So a rewritten block teaches only its edits.
                if lineno not in marks:
                    continue
                for key in course.line_concepts(filename, line):
                    concept = course.CONCEPTS.get(key)
                    if concept and key not in seen:
                        seen.add(key)
                        kind, title, bullets = concept[0], concept[1], concept[2]
                        # An optional 4th element is a short code example, shown
                        # in the monospace sub-line above the bullets.
                        example = concept[3] if len(concept) > 3 else ""
                        out.append(("concept", {"eyebrow": eyebrow_for[kind], "title": title,
                                                "sub": example, "bullets": bullets,
                                                "vis": getattr(course, "VISUALS", {}).get(key)}))
                        # If the concept has a visible spot on the app, follow the
                        # animation with a mockup slide ringing that spot.
                        place = getattr(course, "PLACEMENTS", {}).get(key)
                        if place:
                            out.append(("mockup", {"eyebrow": "On the page", "title": title,
                                                   "region": place[0], "cap": place[1]}))
                note = notes[i] if notes and i < len(notes) else ""
                if note is None:
                    continue  # a line deliberately folded into its neighbour
                # Show the block built up so far, with THIS line highlighted, so
                # the student sees exactly where the new line goes.
                out.append(("linectx", {"file": filename, "start": start,
                                        "lines": lines[:i + 1], "hi": i,
                                        "note": note, "marks": marks}))
                shown += 1
            emit_deletes(start + len(lines), lines[:])   # a line removed from the end
            # One slide with the whole chunk together, once it is more than a
            # single line, so they can check theirs matches before moving on.
            if shown > 1:
                whole = lines[:]
                while whole and not whole[-1].strip():
                    whole.pop()
                out.append(("block", {"file": filename, "start": start,
                                      "lines": whole, "marks": marks}))
            # After the chunk is whole, a quick multiple-choice check: a few
            # questions on the code just typed, plus one or two from earlier weeks
            # to keep old ideas fresh. Each is a question slide then a reveal
            # slide. Authored per block in course.QUIZZES.
            for quiz in getattr(course, "QUIZZES", {}).get((week["n"], filename, block), []):
                q = balance_quiz(quiz)
                out.append(("quiz", {"quiz": q}))
                out.append(("quizanswer", {"quiz": q}))
    return out


def _code_table_html(start, lines, hi, marks):
    """A block of code with line numbers. If hi is set, that row is highlighted
    and the rest are dimmed context; if hi is None, new lines are green."""
    rows = []
    for j, line in enumerate(lines):
        lineno = start + j
        if hi is not None:
            cls = "hi" if j == hi else "ctx"
        else:
            cls = "new" if lineno in marks else ""
        attr = ' class="%s"' % cls if cls else ""
        rows.append('<tr%s><td class="ln">%d</td><td>%s</td></tr>'
                    % (attr, lineno, esc(line) or "&nbsp;"))
    return '<table class="code">%s</table>' % "".join(rows)


# ---------------------------------------------------------------------------
# Visual metaphors for concept slides.
#
# A concept slide replaces its text bullets with a small animated picture that
# SHOWS what the code does (a box sliding onto a list, a token flowing through a
# function, a request bouncing to a server and back). Each concept names one
# "kind" of metaphor in course.VISUALS and fills in a couple of labels; the code
# below turns that spec into a self-contained inline <svg> plus its own scoped
# <style> (concrete @keyframes, uniquely named so slides never clash). No JS and
# no external libraries - it loops forever on CSS alone, so whenever a student
# lands on the slide they see the motion play within a few seconds.
#
# The .pptx cannot animate, so there it falls back to the concept's text bullets
# (still authored in CONCEPTS). Both decks stay one slide, so the drift guard is
# happy.
# ---------------------------------------------------------------------------
_VIS_SEQ = [0]
_VW, _VH = 680, 210          # viewBox
_BW, _BH, _GAP, _PITCH, _ROWY = 88, 64, 16, 104, 73


def _vbox(x, label, cls=""):
    """A list/array cell: rounded box with centred monospace text."""
    return ('<g class="vb {cls}"><rect class="vb-rect" x="{x}" y="{y}" width="{w}" '
            'height="{h}" rx="13"/><text class="vb-txt" x="{tx}" y="{ty}">{t}</text></g>').format(
                cls=cls, x=x, y=_ROWY, w=_BW, h=_BH, tx=x + _BW // 2, ty=_ROWY + _BH // 2,
                t=esc(label))


def _vtile(x, y, w, h, label, cls=""):
    """A small labelled tile (function input/output, network node, string)."""
    return ('<g class="vtile {cls}"><rect x="{x}" y="{y}" width="{w}" height="{h}" rx="10"/>'
            '<text x="{tx}" y="{ty}">{t}</text></g>').format(
                cls=cls, x=x, y=y, w=w, h=h, tx=x + w // 2, ty=y + h // 2, t=esc(label))


def concept_visual(vis):
    """An animated inline SVG for one concept metaphor. Returns HTML."""
    _VIS_SEQ[0] += 1
    cid = "v%d" % _VIS_SEQ[0]
    k = vis["kind"]
    css, body = [], []

    def kf(name, frames):
        css.append("@keyframes %s{%s}" % (name, frames))

    def anim(sel, name, dur, ease="ease-in-out"):
        css.append(".%s{animation:%s %ss %s infinite}" % (sel, name, dur, ease))

    if k in ("arr-add", "arr-remove"):
        items = vis["items"]
        end = vis.get("end", "right")
        n = len(items)
        if k == "arr-add":
            slots = n + 1
            totalw = slots * _BW + (slots - 1) * _GAP
            startx = (_VW - totalw) // 2
            if end == "right":
                xs = [startx + i * _PITCH for i in range(n)]
                fin = startx + n * _PITCH
                off = _VW + 140 - fin
            else:
                xs = [startx + (i + 1) * _PITCH for i in range(n)]
                fin = startx
                off = -(fin + _BW + 120)
            for i, it in enumerate(items):
                body.append(_vbox(xs[i], it))
            body.append('<g class="%s-in">%s</g>' % (cid, _vbox(fin, vis.get("incoming", "new"), "vb-in")))
            kf(cid + "k", "0%%{transform:translateX(%dpx);opacity:0}12%%{opacity:1}"
               "42%%{transform:translateX(0);opacity:1}80%%{transform:translateX(0);opacity:1}"
               "90%%{opacity:0}100%%{transform:translateX(%dpx);opacity:0}" % (off, off))
            anim(cid + "-in", cid + "k", "3.4", "cubic-bezier(.5,0,.2,1)")
        else:
            totalw = n * _BW + (n - 1) * _GAP
            startx = (_VW - totalw) // 2
            xs = [startx + i * _PITCH for i in range(n)]
            li = n - 1 if end == "right" else 0
            off = 240 if end == "right" else -240
            for i, it in enumerate(items):
                if i == li:
                    body.append('<g class="%s-out">%s</g>' % (cid, _vbox(xs[i], it, "vb-in")))
                elif end == "left":
                    body.append('<g class="%s-rest">%s</g>' % (cid, _vbox(xs[i], it)))
                else:
                    body.append(_vbox(xs[i], it))
            kf(cid + "o", "0%%{transform:translateX(0);opacity:1}28%%{transform:translateX(0);opacity:1}"
               "52%%{transform:translateX(%dpx);opacity:0}86%%{opacity:0}"
               "100%%{transform:translateX(0);opacity:1}" % off)
            anim(cid + "-out", cid + "o", "3.4")
            if end == "left":
                kf(cid + "r", "0%%{transform:translateX(0)}28%%{transform:translateX(0)}"
                   "52%%{transform:translateX(-%dpx)}86%%{transform:translateX(-%dpx)}"
                   "100%%{transform:translateX(0)}" % (_PITCH, _PITCH))
                anim(cid + "-rest", cid + "r", "3.4")

    elif k == "loop":
        items = vis["items"]
        n = len(items)
        totalw = n * _BW + (n - 1) * _GAP
        startx = (_VW - totalw) // 2
        body.append('<rect class="%s-hi vb-hi" x="%d" y="%d" width="%d" height="%d" rx="14"/>'
                     % (cid, startx, _ROWY - 6, _BW, _BH + 12))
        for i, it in enumerate(items):
            body.append(_vbox(startx + i * _PITCH, it))
        seg = 100 // n
        frames = []
        for i in range(n):
            a, hold = i * seg, i * seg + max(1, seg // 2)
            frames.append("%d%%{transform:translateX(%dpx)}" % (a, i * _PITCH))
            frames.append("%d%%{transform:translateX(%dpx)}" % (hold, i * _PITCH))
        frames.append("100%{transform:translateX(0)}")
        kf(cid + "k", "".join(frames))
        anim(cid + "-hi", cid + "k", "%.1f" % max(2.6, n * 0.9))

    elif k == "machine":
        mw, mh = 190, 96
        mx, my = (_VW - mw) // 2, (_VH - mh) // 2
        tw, th = 92, 54
        ty = (_VH - th) // 2
        body.append('<rect class="vmachine" x="%d" y="%d" width="%d" height="%d" rx="16"/>'
                     % (mx, my, mw, mh))
        body.append('<text class="vmlabel" x="%d" y="%d">%s</text>'
                     % (mx + mw // 2, my + mh // 2, esc(vis.get("label", "f()"))))
        body.append('<g class="%s-in">%s</g>' % (cid, _vtile(mx - tw - 66, ty, tw, th, vis.get("in", "in"), "vt-in")))
        body.append('<g class="%s-out">%s</g>' % (cid, _vtile(mx + mw + 66, ty, tw, th, vis.get("out", "out"), "vt-out")))
        d = tw + 66
        kf(cid + "i", "0%%{transform:translateX(0);opacity:0}9%%{opacity:1}"
           "40%%{transform:translateX(%dpx);opacity:1}48%%{transform:translateX(%dpx);opacity:0}"
           "100%%{opacity:0}" % (d, d + 20))
        anim(cid + "-in", cid + "i", "3.6")
        kf(cid + "o", "0%%{transform:translateX(-%dpx);opacity:0}52%%{transform:translateX(-%dpx);opacity:0}"
           "62%%{transform:translateX(0);opacity:1}92%%{transform:translateX(0);opacity:1}"
           "100%%{transform:translateX(-%dpx);opacity:0}" % (d, d, d))
        anim(cid + "-out", cid + "o", "3.6")

    elif k == "network":
        nh = 64
        lw, rw = 150, 176
        lx, rx = 66, _VW - 66 - rw
        cy = _ROWY + nh // 2
        body.append(_vtile(lx, _ROWY, lw, nh, vis.get("from", "you")))
        body.append(_vtile(rx, _ROWY, rw, nh, vis.get("to", "server"), "vt-cloud"))
        dist = rx - (lx + lw) - 16
        body.append('<circle class="%s-req" cx="%d" cy="%d" r="12"/>' % (cid, lx + lw + 8, cy))
        kf(cid + "q", "0%%{transform:translateX(0);opacity:0}6%%{opacity:1}"
           "40%%{transform:translateX(%dpx);opacity:1}46%%{transform:translateX(%dpx);opacity:0}"
           "100%%{opacity:0}" % (dist, dist))
        css.append(".%s-req{fill:#ffd633}" % cid)
        anim(cid + "-req", cid + "q", "3.6")
        body.append('<circle class="%s-rep" cx="%d" cy="%d" r="12"/>' % (cid, rx - 8, cy))
        kf(cid + "p", "0%%{opacity:0}54%%{transform:translateX(0);opacity:0}60%%{opacity:1}"
           "94%%{transform:translateX(-%dpx);opacity:1}100%%{transform:translateX(-%dpx);opacity:0}" % (dist, dist))
        css.append(".%s-rep{fill:#7fe0a0}" % cid)
        anim(cid + "-rep", cid + "p", "3.6")

    elif k == "glue":
        tw, th, y1 = 132, 58, 36
        ax, bx = _VW // 2 - tw - 42, _VW // 2 + 42
        body.append(_vtile(ax, y1, tw, th, vis.get("a", '"Hi "')))
        body.append('<text class="vplus" x="%d" y="%d">+</text>' % (_VW // 2, y1 + th // 2))
        body.append(_vtile(bx, y1, tw, th, vis.get("b", "name")))
        rw = 260
        body.append('<text class="varrow" x="%d" y="122">&#8595;</text>' % (_VW // 2))
        body.append('<g class="%s-r">%s</g>' % (cid, _vtile((_VW - rw) // 2, 138, rw, th, vis.get("out", '"Hi Sam"'), "vt-out")))
        kf(cid + "r", "0%{opacity:0}38%{opacity:0}52%{opacity:1}92%{opacity:1}100%{opacity:0}")
        anim(cid + "-r", cid + "r", "3.2")

    elif k == "swap":
        bw, bh = 240, 90
        bx, by = (_VW - bw) // 2, (_VH - bh) // 2
        tx, ty = bx + bw // 2, by + bh // 2
        body.append('<g class="%s-off"><rect class="vswap-off" x="%d" y="%d" width="%d" height="%d" rx="15"/>'
                     '<text class="vb-txt" x="%d" y="%d">%s</text></g>'
                     % (cid, bx, by, bw, bh, tx, ty, esc(vis.get("off", "off"))))
        body.append('<g class="%s-on"><rect class="vswap-on" x="%d" y="%d" width="%d" height="%d" rx="15"/>'
                     '<text class="vb-txt" x="%d" y="%d">%s</text></g>'
                     % (cid, bx, by, bw, bh, tx, ty, esc(vis.get("on", "on"))))
        kf(cid + "f", "0%{opacity:1}44%{opacity:1}56%{opacity:0}94%{opacity:0}100%{opacity:1}")
        kf(cid + "n", "0%{opacity:0}44%{opacity:0}56%{opacity:1}94%{opacity:1}100%{opacity:0}")
        anim(cid + "-off", cid + "f", "3.0")
        anim(cid + "-on", cid + "n", "3.0")

    elif k == "boxmodel":
        layer = vis.get("layer", "padding")
        w, h = 380, 156
        x, y = (_VW - w) // 2, (_VH - h) // 2
        rects = {
            "margin": (x, y, w, h),
            "border": (x + 26, y + 26, w - 52, h - 52),
            "padding": (x + 50, y + 50, w - 100, h - 100),
        }
        colours = {"margin": "#ffd633", "border": "#c98be0", "padding": "#7fe0a0"}
        for name in ("margin", "border", "padding"):
            rx, ry, rw2, rh2 = rects[name]
            dash = ' stroke-dasharray="7 6"' if name == "margin" else ""
            cls = ' class="%s-pulse"' % cid if name == layer else ""
            body.append('<rect%s x="%d" y="%d" width="%d" height="%d" rx="8" fill="none" '
                        'stroke="%s" stroke-width="3"%s/>' % (cls, rx, ry, rw2, rh2, colours[name], dash))
        cxr = rects["padding"]
        body.append('<rect x="%d" y="%d" width="%d" height="%d" rx="6" fill="#132833"/>'
                     % (cxr[0] + 20, cxr[1] + 18, cxr[2] - 40, cxr[3] - 36))
        body.append('<text class="vb-txt" x="%d" y="%d" style="font-size:22px">content</text>'
                     % (_VW // 2, _VH // 2))
        kf(cid + "k", "0%{stroke-width:3;opacity:.55}50%{stroke-width:7;opacity:1}100%{stroke-width:3;opacity:.55}")
        anim(cid + "-pulse", cid + "k", "1.8")

    elif k == "box":                                   # a variable = a labelled box
        bw, bh, by = 220, 90, 106
        bx = (_VW - bw) // 2
        body.append('<text class="vmlabel" x="%d" y="66" style="font-size:24px">%s</text>'
                     % (_VW // 2, esc(vis.get("name", "value"))))
        body.append('<rect class="vmachine" x="%d" y="%d" width="%d" height="%d" rx="14"/>' % (bx, by, bw, bh))
        tw, th = 130, 52
        body.append('<g class="%s-v">%s</g>' % (cid, _vtile((_VW - tw) // 2, by + (bh - th) // 2, tw, th, vis.get("value", "0"), "vt-in")))
        kf(cid + "k", "0%{transform:translateY(-92px);opacity:0}18%{opacity:1}44%{transform:translateY(6px);opacity:1}"
           "54%{transform:translateY(0)}86%{transform:translateY(0);opacity:1}94%{opacity:0}100%{transform:translateY(-92px);opacity:0}")
        anim(cid + "-v", cid + "k", "3.2")

    elif k == "pick":                                  # a pointer picking one cell
        items, at = vis["items"], vis.get("at", 0)
        n = len(items)
        totalw = n * _BW + (n - 1) * _GAP
        startx = (_VW - totalw) // 2
        tx = startx + at * _PITCH
        for i, it in enumerate(items):
            body.append(_vbox(startx + i * _PITCH, it, "vb-pick" if i == at else ""))
        body.append('<g class="%s-pt"><text class="vmlabel" x="%d" y="34" style="font-size:22px">%s</text>'
                     '<text class="varrow" x="%d" y="60" style="font-size:34px">&#8595;</text></g>'
                     % (cid, tx + _BW // 2, esc(vis.get("label", "[0]")), tx + _BW // 2))
        kf(cid + "k", "0%{transform:translateY(0)}50%{transform:translateY(-9px)}100%{transform:translateY(0)}")
        anim(cid + "-pt", cid + "k", "1.5")

    elif k == "fork":                                  # a check that branches two ways
        cxc = _VW // 2
        body.append('<rect class="vmachine" x="%d" y="16" width="180" height="58" rx="12"/>' % (cxc - 90))
        body.append('<text class="vmlabel" x="%d" y="45" style="font-size:21px">%s</text>' % (cxc, esc(vis.get("cond", "check"))))
        tw, th, ly = 156, 58, 140
        lx, rx = cxc - tw - 34, cxc + 34
        body.append('<path d="M%d,74 L%d,%d" stroke="#33586a" stroke-width="2.5" fill="none"/>' % (cxc, lx + tw // 2, ly))
        body.append('<path d="M%d,74 L%d,%d" stroke="#33586a" stroke-width="2.5" fill="none"/>' % (cxc, rx + tw // 2, ly))
        body.append(_vtile(lx, ly, tw, th, vis.get("yes", "true"), "vt-out"))
        body.append(_vtile(rx, ly, tw, th, vis.get("no", "false"), "vt-no"))
        body.append('<circle class="%s-tok" cx="%d" cy="74" r="10"/>' % (cid, cxc))
        dx, dy = lx + tw // 2 - cxc, ly - 74
        kf(cid + "k", "0%%{transform:translate(0,0);opacity:0}12%%{opacity:1}58%%{transform:translate(%dpx,%dpx);opacity:1}"
           "78%%{transform:translate(%dpx,%dpx);opacity:0}100%%{opacity:0}" % (dx, dy, dx, dy))
        css.append(".%s-tok{fill:#ffd633}" % cid)
        anim(cid + "-tok", cid + "k", "3.0")

    elif k == "event":                                 # a click sparks code to run
        bx, by, bw, bh = 84, _ROWY, 156, 64
        body.append(_vtile(bx, by, bw, bh, vis.get("btn", "button"), "vt-btn"))
        body.append('<circle class="%s-rip" cx="%d" cy="%d" r="16" fill="none" stroke="#ffd633" stroke-width="3"/>'
                     % (cid, bx + bw // 2, by + bh // 2))
        kf(cid + "r", "0%{r:14px;opacity:.9}70%{r:54px;opacity:0}100%{r:54px;opacity:0}")
        anim(cid + "-rip", cid + "r", "2.6")
        cx2 = _VW - 84 - 200
        body.append('<g class="%s-code">%s</g>' % (cid, _vtile(cx2, by, 200, bh, vis.get("action", "run code"), "vt-out")))
        body.append('<circle class="%s-sp" cx="%d" cy="%d" r="8"/>' % (cid, bx + bw + 8, by + bh // 2))
        dist = cx2 - (bx + bw) - 16
        kf(cid + "s", "0%%{transform:translateX(0);opacity:0}20%%{opacity:1}54%%{transform:translateX(%dpx);opacity:1}"
           "62%%{opacity:0}100%%{opacity:0}" % dist)
        css.append(".%s-sp{fill:#ffd633}" % cid)
        anim(cid + "-sp", cid + "s", "2.6")
        kf(cid + "f", "0%{opacity:.55}56%{opacity:.55}64%{opacity:1}90%{opacity:1}100%{opacity:.55}")
        anim(cid + "-code", cid + "f", "2.6")

    elif k == "dom":                                   # a child node nests in a parent
        pw, ph = 320, 128
        px, py = (_VW - pw) // 2, (_VH - ph) // 2
        body.append('<rect class="vmachine" x="%d" y="%d" width="%d" height="%d" rx="14"/>' % (px, py, pw, ph))
        body.append('<text class="vmlabel" x="%d" y="%d" style="font-size:20px">%s</text>'
                     % (px + pw // 2, py + 22, esc(vis.get("parent", "parent"))))
        tw, th = 168, 52
        tx, ty = px + (pw - tw) // 2, py + ph - th - 16
        body.append('<g class="%s-c">%s</g>' % (cid, _vtile(tx, ty, tw, th, vis.get("child", "child"), "vt-out")))
        if vis.get("mode", "add") == "add":
            kf(cid + "k", "0%{transform:translateY(-128px);opacity:0}20%{opacity:1}48%{transform:translateY(0);opacity:1}"
               "86%{transform:translateY(0);opacity:1}94%{opacity:0}100%{transform:translateY(-128px);opacity:0}")
        else:
            kf(cid + "k", "0%{transform:translateY(0);opacity:1}30%{transform:translateY(0);opacity:1}"
               "56%{transform:translateY(-128px);opacity:0}100%{transform:translateY(-128px);opacity:0}")
        anim(cid + "-c", cid + "k", "3.2")

    elif k == "tag":                                   # opening + closing tags hug content
        name = vis.get("tag", "p")
        if vis.get("selfclose"):
            lbl = vis.get("open", "<%s />" % name)
            body.append(_vtile((_VW - 300) // 2, _ROWY, 300, 64, lbl, "vt-tag"))
        else:
            openl = vis.get("open", "<%s>" % name)
            closel = vis.get("close", "</%s>" % name)
            ow, cw, conw, g = 130, 130, 230, 12
            total = ow + conw + cw + 2 * g
            sx = (_VW - total) // 2
            body.append('<g class="%s-o">%s</g>' % (cid, _vtile(sx, _ROWY, ow, 64, openl, "vt-tag")))
            body.append(_vtile(sx + ow + g, _ROWY, conw, 64, vis.get("content", "text")))
            body.append('<g class="%s-c">%s</g>' % (cid, _vtile(sx + ow + g + conw + g, _ROWY, cw, 64, closel, "vt-tag")))
            kf(cid + "o", "0%{transform:translateX(0)}50%{transform:translateX(11px)}100%{transform:translateX(0)}")
            kf(cid + "c", "0%{transform:translateX(0)}50%{transform:translateX(-11px)}100%{transform:translateX(0)}")
            anim(cid + "-o", cid + "o", "2.2")
            anim(cid + "-c", cid + "c", "2.2")

    elif k == "swatch":                                # a colour that shifts
        target = vis.get("target", "box")
        a, b = vis.get("a", "#1f6feb"), vis.get("b", "#7c5cff")
        bw, bh = 280, 108
        bx, by = (_VW - bw) // 2, (_VH - bh) // 2 - 8
        if target == "text":
            body.append('<rect x="%d" y="%d" width="%d" height="%d" rx="14" fill="#132833" stroke="#33586a" stroke-width="2.5"/>' % (bx, by, bw, bh))
            body.append('<text class="%s-sw" x="%d" y="%d" style="font:800 40px Rubik,sans-serif;text-anchor:middle;dominant-baseline:central">Text</text>' % (cid, _VW // 2, by + bh // 2))
            kf(cid + "k", "0%%{fill:%s}50%%{fill:%s}100%%{fill:%s}" % (a, b, a))
        elif target == "border":
            body.append('<rect class="%s-sw" x="%d" y="%d" width="%d" height="%d" rx="14" fill="#132833" stroke="%s" stroke-width="7"/>' % (cid, bx, by, bw, bh, a))
            kf(cid + "k", "0%%{stroke:%s}50%%{stroke:%s}100%%{stroke:%s}" % (a, b, a))
        else:
            body.append('<rect class="%s-sw" x="%d" y="%d" width="%d" height="%d" rx="14" stroke="#33586a" stroke-width="2" fill="%s"/>' % (cid, bx, by, bw, bh, a))
            kf(cid + "k", "0%%{fill:%s}50%%{fill:%s}100%%{fill:%s}" % (a, b, a))
        css.append(".%s-sw{animation:%sk 3.0s ease-in-out infinite}" % (cid, cid))
        body.append('<text x="%d" y="%d" style="font:700 20px Consolas,monospace;text-anchor:middle;fill:#8fb3bd">%s</text>'
                     % (_VW // 2, by + bh + 26, esc(vis.get("label", b))))

    elif k == "resize":                                # a box that grows on one axis
        bw, bh = 128, 128
        body.append('<rect class="%s-rz" x="%d" y="%d" width="%d" height="%d" rx="12" fill="#0f2c38" stroke="#01aefd" stroke-width="3"/>'
                     % (cid, (_VW - bw) // 2, (_VH - bh) // 2, bw, bh))
        if vis.get("axis", "w") == "w":
            kf(cid + "k", "0%{transform:scaleX(.45)}50%{transform:scaleX(1.7)}100%{transform:scaleX(.45)}")
        else:
            kf(cid + "k", "0%{transform:scaleY(.45)}50%{transform:scaleY(1.5)}100%{transform:scaleY(.45)}")
        css.append(".%s-rz{transform-box:fill-box;transform-origin:center;animation:%sk 2.6s ease-in-out infinite}" % (cid, cid))

    elif k == "round":                                 # corners rounding off
        bw, bh = 150, 150
        body.append('<rect class="%s-r" x="%d" y="%d" width="%d" height="%d" fill="#0f2c38" stroke="#01aefd" stroke-width="3"/>'
                     % (cid, (_VW - bw) // 2, (_VH - bh) // 2, bw, bh))
        kf(cid + "k", "0%{rx:2px;ry:2px}50%{rx:46px;ry:46px}100%{rx:2px;ry:2px}")
        css.append(".%s-r{animation:%sk 2.6s ease-in-out infinite}" % (cid, cid))

    elif k == "textsize":                              # text growing
        body.append('<text class="%s-t" x="%d" y="%d" style="fill:#01aefd;text-anchor:middle;dominant-baseline:central;font-weight:800;font-family:Rubik,sans-serif">Aa</text>'
                     % (cid, _VW // 2, _VH // 2))
        kf(cid + "k", "0%{font-size:32px}50%{font-size:96px}100%{font-size:32px}")
        css.append(".%s-t{animation:%sk 2.6s ease-in-out infinite}" % (cid, cid))

    elif k == "linegap":                               # line spacing opening up
        for tag2, lbl, base in (("a", "line one", -24), ("b", "line two", 24)):
            body.append('<g class="%s-%s"><text x="%d" y="%d" style="fill:#e6f2f6;text-anchor:middle;dominant-baseline:central;font:700 26px Rubik,sans-serif">%s</text></g>'
                         % (cid, tag2, _VW // 2, _VH // 2 + base, lbl))
        kf(cid + "a", "0%{transform:translateY(0)}50%{transform:translateY(-22px)}100%{transform:translateY(0)}")
        kf(cid + "b", "0%{transform:translateY(0)}50%{transform:translateY(22px)}100%{transform:translateY(0)}")
        anim(cid + "-a", cid + "a", "2.6")
        anim(cid + "-b", cid + "b", "2.6")

    elif k == "motion":                                # something moving on its own
        cy = _VH // 2
        body.append('<line x1="72" y1="%d" x2="%d" y2="%d" stroke="#33586a" stroke-width="2" stroke-dasharray="6 6"/>' % (cy, _VW - 72, cy))
        body.append('<circle class="%s-m" cx="90" cy="%d" r="18" fill="#01aefd"/>' % (cid, cy))
        kf(cid + "k", "0%%{transform:translateX(0)}50%%{transform:translateX(%dpx)}100%%{transform:translateX(0)}" % (_VW - 180))
        anim(cid + "-m", cid + "k", "2.4")

    elif k == "scroll":                                # content scrolling in a viewport
        vw2, vh2 = 300, 140
        vx, vy = (_VW - vw2) // 2, (_VH - vh2) // 2
        body.append('<clipPath id="%s-clip"><rect x="%d" y="%d" width="%d" height="%d" rx="10"/></clipPath>' % (cid, vx, vy, vw2, vh2))
        body.append('<rect x="%d" y="%d" width="%d" height="%d" rx="10" fill="#0f1b21" stroke="#33586a" stroke-width="2.5"/>' % (vx, vy, vw2, vh2))
        rowsn, lh = 8, 30
        inner = "".join('<rect x="%d" y="%d" width="%d" height="16" rx="4" fill="#2a4a57"/>' % (vx + 18, vy + 14 + i * lh, vw2 - 70) for i in range(rowsn))
        body.append('<g clip-path="url(#%s-clip)"><g class="%s-sc">%s</g></g>' % (cid, cid, inner))
        dist = rowsn * lh + 14 - vh2
        kf(cid + "k", "0%%{transform:translateY(0)}45%%{transform:translateY(-%dpx)}55%%{transform:translateY(-%dpx)}100%%{transform:translateY(0)}" % (dist, dist))
        anim(cid + "-sc", cid + "k", "4.0", "linear")
        body.append('<rect class="%s-th" x="%d" y="%d" width="6" height="40" rx="3" fill="#01aefd"/>' % (cid, vx + vw2 - 12, vy + 8))
        kf(cid + "t", "0%%{transform:translateY(0)}45%%{transform:translateY(%dpx)}55%%{transform:translateY(%dpx)}100%%{transform:translateY(0)}" % (vh2 - 56, vh2 - 56))
        anim(cid + "-th", cid + "t", "4.0", "linear")

    elif k == "card":                                  # an object as a key:value card
        rows = vis.get("rows", [("role", "user"), ("text", "hi")])
        n = len(rows)
        cw, rh = 380, 46
        ch = 30 + n * rh
        cx, cy = (_VW - cw) // 2, (_VH - ch) // 2
        body.append('<rect class="vmachine" x="%d" y="%d" width="%d" height="%d" rx="14"/>' % (cx, cy, cw, ch))
        for i, (kk, vv) in enumerate(rows):
            ry = cy + 22 + i * rh
            body.append('<text x="%d" y="%d" style="fill:#7fd8ff;font:700 22px Consolas,monospace;dominant-baseline:central">%s:</text>' % (cx + 26, ry, esc(kk)))
            body.append('<text x="%d" y="%d" style="fill:#e6f2f6;font:700 22px Consolas,monospace;text-anchor:end;dominant-baseline:central">%s</text>' % (cx + cw - 26, ry, esc(vv)))
        hy = cy + 22 - (rh - 8) // 2
        body.append('<rect class="%s-hl vb-hi" x="%d" y="%d" width="%d" height="%d" rx="8"/>' % (cid, cx + 10, hy, cw - 20, rh - 8))
        frames = []
        for i in range(n):
            a, hold = i * (100 // n), i * (100 // n) + (100 // n) // 2
            frames.append("%d%%{transform:translateY(%dpx)}" % (a, i * rh))
            frames.append("%d%%{transform:translateY(%dpx)}" % (hold, i * rh))
        frames.append("100%{transform:translateY(0)}")
        kf(cid + "k", "".join(frames))
        css.append(".%s-hl{opacity:.55;animation:%sk %.1fs ease-in-out infinite}" % (cid, cid, max(2.4, n * 0.9)))

    elif k == "tiles":                                 # labelled parts of a line, swept
        parts = vis["parts"]
        n = len(parts)
        ws = [max(84, 28 + len(str(p)) * 15) for p in parts]
        g = 14
        total = sum(ws) + g * (n - 1)
        x = (_VW - total) // 2
        xs = []
        for w in ws:
            xs.append(x)
            x += w + g
        for i, p in enumerate(parts):
            body.append(_vtile(xs[i], _ROWY, ws[i], 64, p))
        dur = max(2.4, n * 0.85)
        kf(cid + "k", "0%{opacity:0}12%{opacity:.85}26%{opacity:0}100%{opacity:0}")
        for i in range(n):
            body.append('<rect class="%s-h%d vb-hi" x="%d" y="%d" width="%d" height="76" rx="12"/>' % (cid, i, xs[i], _ROWY - 6, ws[i]))
            css.append(".%s-h%d{opacity:0;animation:%sk %.1fs ease-in-out infinite;animation-delay:%.2fs}" % (cid, i, cid, dur, i * dur / n))

    elif k == "flex":                                  # items settling into a layout
        mode = vis.get("mode", "row")
        items = vis.get("items", ["A", "B", "C"])
        n = len(items)
        tw, th = 104, 64
        rowy = _VH // 2 - th // 2

        def _spread(gap):
            tot = n * tw + (n - 1) * gap
            sx = (_VW - tot) // 2
            return [(sx + i * (tw + gap), rowy) for i in range(n)]

        if mode == "gap":
            start, target = _spread(2), _spread(48)
        elif mode == "column":
            start = _spread(20)
            th = 50
            ch = n * th + (n - 1) * 8
            sy = (_VH - ch) // 2
            target = [((_VW - tw) // 2, sy + i * (th + 8)) for i in range(n)]
        elif mode == "end":
            left = [(40 + i * (tw + 16), rowy) for i in range(n)]
            start = left[:]
            target = left[:]
            target[-1] = (_VW - 40 - tw, rowy)
        elif mode == "wrap":
            per = (n + 1) // 2

            def _rowpos(idxs, y):
                m = len(idxs)
                tot = m * tw + (m - 1) * 16
                sx = (_VW - tot) // 2
                return {idxs[j]: (sx + j * (tw + 16), y) for j in range(m)}
            tpos = {}
            tpos.update(_rowpos(list(range(per)), rowy - 40))
            tpos.update(_rowpos(list(range(per, n)), rowy + 40))
            start = [((_VW - tw) // 2, rowy)] * n
            target = [tpos[i] for i in range(n)]
        else:  # row
            start, target = [((_VW - tw) // 2, rowy)] * n, _spread(20)
        for i, it in enumerate(items):
            s, t = start[i], target[i]
            body.append('<g class="%s-i%d">%s</g>' % (cid, i, _vtile(s[0], s[1], tw, th, it)))
            kf("%si%d" % (cid, i), "0%%{transform:translate(0,0)}42%%{transform:translate(%dpx,%dpx)}"
               "84%%{transform:translate(%dpx,%dpx)}100%%{transform:translate(0,0)}"
               % (t[0] - s[0], t[1] - s[1], t[0] - s[0], t[1] - s[1]))
            anim("%s-i%d" % (cid, i), "%si%d" % (cid, i), "3.6")

    svg = '<svg viewBox="0 0 %d %d" role="img" aria-label="%s">%s</svg>' % (
        _VW, _VH, esc(vis.get("cap", "")), "".join(body))
    return '<div class="vis"><style>%s</style>%s</div>' % ("".join(css), svg)


# ---------------------------------------------------------------------------
# The mockup: a fixed wireframe of the chat app the students are building. After
# a concept's animation slide, a second slide shows this same little page with
# ONE region ringed, so they see where that piece actually lives. The wireframe
# never changes shape slide to slide (so they learn to recognise it) - only the
# highlighted region moves. course.PLACEMENTS maps a concept to a region + line.
# ---------------------------------------------------------------------------
_MOCK_VW, _MOCK_VH = 640, 440


def _mock_regions():
    """name -> (x, y, w, h, corner-radius) in the wireframe's viewBox units."""
    return {
        "page":       (44, 14, 552, 412, 14),
        "header":     (60, 54, 520, 46, 8),
        "title":      (72, 68, 150, 20, 4),
        "log":        (60, 110, 520, 190, 8),
        "aibubble":   (72, 122, 300, 42, 12),
        "userbubble": (268, 176, 300, 42, 12),
        "inputrow":   (60, 344, 520, 50, 10),
        "textbox":    (70, 352, 410, 34, 8),
        "sendbtn":    (490, 352, 82, 34, 8),
        "footer":     (60, 404, 520, 22, 6),
    }


def concept_mockup(region):
    """The wireframe with `region` ringed in a pulsing highlight. Returns HTML."""
    _VIS_SEQ[0] += 1
    cid = "m%d" % _VIS_SEQ[0]
    b = []
    # browser window + title bar
    b.append('<rect x="44" y="14" width="552" height="412" rx="14" fill="#0f1f27" stroke="#33586a" stroke-width="2.5"/>')
    b.append('<path d="M44 30 q0 -16 16 -16 h520 q16 0 16 16 v14 h-552 z" fill="#16303a"/>')
    for cxx in (66, 84, 102):
        b.append('<circle cx="%d" cy="30" r="4" fill="#385c6a"/>' % cxx)
    # header with a title block
    b.append('<rect x="60" y="54" width="520" height="46" rx="8" fill="#16303a"/>')
    b.append('<rect x="72" y="68" width="150" height="20" rx="4" fill="#4a6b78"/>')
    # chat log with three message bubbles (ai left, you right, ai left)
    b.append('<rect x="60" y="110" width="520" height="190" rx="8" fill="#0b171d" stroke="#22424d" stroke-width="1.5"/>')
    b.append('<rect x="72" y="122" width="300" height="42" rx="12" fill="#1b3540"/>')
    b.append('<rect x="86" y="137" width="220" height="12" rx="4" fill="#3a5b68"/>')
    b.append('<rect x="268" y="176" width="300" height="42" rx="12" fill="#123a32"/>')
    b.append('<rect x="282" y="191" width="220" height="12" rx="4" fill="#2f6a58"/>')
    b.append('<rect x="72" y="230" width="330" height="42" rx="12" fill="#1b3540"/>')
    b.append('<rect x="86" y="245" width="250" height="12" rx="4" fill="#3a5b68"/>')
    # input row: a wide text box and a Send button
    b.append('<rect x="70" y="352" width="410" height="34" rx="8" fill="#16303a" stroke="#33586a" stroke-width="1.5"/>')
    b.append('<rect x="490" y="352" width="82" height="34" rx="8" fill="#01aefd"/>')
    b.append('<text x="531" y="369" style="fill:#04222f;font:700 15px Rubik,system-ui,sans-serif;text-anchor:middle;dominant-baseline:central">Send</text>')
    # footer credits line
    b.append('<rect x="72" y="410" width="170" height="10" rx="3" fill="#2a4550"/>')
    # the highlight ring on the target region
    x, y, w, h, rx = _mock_regions().get(region, _mock_regions()["page"])
    b.append('<rect class="%s-hl" x="%d" y="%d" width="%d" height="%d" rx="%d" fill="none" stroke="#ffd633" stroke-width="4"/>'
             % (cid, x - 4, y - 4, w + 8, h + 8, rx + 4))
    css = ("@keyframes %sk{0%%,100%%{opacity:.4;stroke-width:3}50%%{opacity:1;stroke-width:5.5}}"
           ".%s-hl{animation:%sk 1.5s ease-in-out infinite}" % (cid, cid, cid))
    svg = '<svg viewBox="0 0 %d %d" role="img">%s</svg>' % (_MOCK_VW, _MOCK_VH, "".join(b))
    return '<div class="vis mock"><style>%s</style>%s</div>' % (css, svg)


def deck_render_html(desc):
    """One deck slide as HTML, for any descriptor slide_plan() produces."""
    kind, d = desc
    if kind == "concept":
        eyebrow = f'<p class="eyebrow">{esc(d["eyebrow"])}</p>' if d.get("eyebrow") else ""
        sub = ""
        if d.get("sub"):
            code_ish = any(mark in d["sub"] for mark in ("(", "{", "[", "=", ".", "/", ":", ";", "<"))
            sub = '<p class="sub{0}">{1}</p>'.format(" mono" if code_ish else "", esc(d["sub"]))
        # A concept with a visual metaphor shows the animation and one caption
        # in place of the text bullets; without one it keeps the bullet list.
        if d.get("vis"):
            vis = d["vis"]
            cap = '<p class="vis-cap">{0}</p>'.format(esc(vis["cap"])) if vis.get("cap") else ""
            return ('<section class="slide concept-vis">{0}<h2>{1}</h2>{2}{3}{4}</section>').format(
                eyebrow, esc(d["title"]), sub, concept_visual(vis), cap)
        pts = "".join("<li>{0}</li>".format(esc(pt)) for pt in d.get("bullets", []))
        pts = '<ul class="pts">{0}</ul>'.format(pts) if pts else ""
        return '<section class="slide">{0}<h2>{1}</h2>{2}{3}</section>'.format(
            eyebrow, esc(d["title"]), sub, pts)
    if kind == "mockup":
        cap = '<p class="vis-cap">{0}</p>'.format(esc(d["cap"])) if d.get("cap") else ""
        return ('<section class="slide concept-vis"><p class="eyebrow">{0}</p>'
                '<h2>{1}</h2>{2}{3}</section>').format(
                    esc(d.get("eyebrow", "On the page")), esc(d["title"]),
                    concept_mockup(d["region"]), cap)
    if kind == "checkpoint":
        page = checkpoint_page(d["week"]).replace("</script", "<\\/script")
        cid = "cp%d" % d["week"]
        return ('<section class="slide checkpoint">'
                '<p class="eyebrow">Checkpoint &middot; run it</p>'
                '<h2>{0}</h2><p class="cp-say">{1}</p>'
                '<button class="cp-run" data-cp="{2}">&#9654; Show the output</button>'
                '<div class="cp-out" id="out-{2}"></div>'
                '<script type="text/plain" id="src-{2}">{3}</script>'
                '</section>').format(esc(d["title"]), esc(d["say"]), cid, page)
    if kind in ("quiz", "quizanswer"):
        q = d["quiz"]
        answered = kind == "quizanswer"
        opts = "".join(
            '<li class="{0}">{1}</li>'.format(
                "right" if (answered and i == q["answer"]) else "", esc(opt))
            for i, opt in enumerate(q["options"]))
        tail = ('<p class="quiz-why">{0}</p>'.format(esc(q["why"]))
                if answered else
                '<p class="quiz-hint">Pick one &mdash; the answer is on the next slide.</p>')
        return ('<section class="slide quiz{0}"><p class="eyebrow">{1}</p>'
                '<h2>{2}</h2><ol class="opts">{3}</ol>{4}</section>').format(
                    " answered" if answered else "", "Answer" if answered else "Quick check",
                    esc(q["q"]), opts, tail)
    if kind == "chunk":
        return deck_code_slide(d["file"], d["start"], d["lines"], d["marks"], d["gone"],
                               d["part"], d["parts"])
    if kind == "block":
        return ('<section class="slide dark line whole">'
                '<p class="filebar">All together in {0}</p>{1}'
                '<p class="linenote">That is the whole piece. Check yours looks the '
                'same before moving on.</p></section>').format(
                    esc(d["file"]), _code_table_html(d["start"], d["lines"], None, d["marks"]))
    if kind == "delete":
        rows = ['<tr class="ctx"><td class="ln">{0}</td><td>{1}</td></tr>'.format(
                    d["start"] + j, esc(line) or "&nbsp;") for j, line in enumerate(d["lines"])]
        # Show the real line number of the line to remove (it sits just below the
        # context), not a dash - a student cannot find "the struck line" without
        # the number they see in their own editor.
        rows.append('<tr class="gone hi"><td class="ln">{0}</td><td>{1}</td></tr>'.format(
            d["at"], esc(d["dropped"]) or "&nbsp;"))
        label = "replace a line" if d.get("replacing") else "delete a line"
        return ('<section class="slide dark line del">'
                '<p class="filebar">In {0} <span class="part">{1}</span></p>'
                '<table class="code">{2}</table><p class="linenote">{3}</p></section>').format(
                    esc(d["file"]), label, "".join(rows), esc(d["note"]))
    # a single line in context, with its explanation
    return ('<section class="slide dark line">'
            '<p class="filebar">Type this into {0} <span class="part">line {1}</span></p>'
            '{2}<p class="linenote">{3}</p></section>').format(
                esc(d["file"]), d["start"] + d["hi"],
                _code_table_html(d["start"], d["lines"], d["hi"], d["marks"]),
                esc(d["note"]))


def build_slides():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("  slides SKIPPED - run: pip install python-pptx")
        return 0
    os.makedirs(SLIDES_DIR, exist_ok=True)
    BRAND = RGBColor(0x01, 0xAE, 0xFD)
    INK = RGBColor(0x1F, 0x2A, 0x37)
    PAPER = RGBColor(0xDC, 0xEC, 0xEF)
    NEW = RGBColor(0x7F, 0xE0, 0xA0)
    GUTTER = RGBColor(0x6D, 0x8C, 0x97)
    DROP = RGBColor(0xF3, 0x92, 0x8B)   # a line taken out this week
    DARK = RGBColor(0x10, 0x21, 0x27)

    def concept_slide(deck, spec, eyebrow=""):
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = spec["title"]
        run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.color.rgb = Pt(38), True, INK
        if eyebrow:
            eb = slide.shapes.add_textbox(Inches(0.95), Inches(0.55), Inches(6), Inches(0.5))
            eb.text_frame.text = eyebrow.upper()
            er = eb.text_frame.paragraphs[0].runs[0]
            er.font.size, er.font.bold, er.font.color.rgb = Pt(15), True, BRAND
        body = ([spec["sub"]] if spec.get("sub") else []) + spec.get("bullets", [])
        if not body:
            return
        box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
        frame = box.text_frame
        frame.word_wrap = True
        for i, line in enumerate(body):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            is_bullet = line in spec.get("bullets", [])
            para.text = ("• " + line) if is_bullet else line
            para.runs[0].font.size = Pt(24 if is_bullet else 28)
            para.runs[0].font.color.rgb = INK if is_bullet else BRAND
            para.space_after = Pt(14)

    def mockup_slide(deck, spec):
        """A static wireframe of the app with one region ringed - the .pptx
        twin of concept_mockup, so Google Slides shows the placement too."""
        YELLOW = RGBColor(0xFF, 0xD6, 0x33)
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(0xF4, 0xF8, 0xFB)
        eb = slide.shapes.add_textbox(Inches(0.9), Inches(0.4), Inches(8), Inches(0.4))
        eb.text_frame.text = "ON THE PAGE"
        r = eb.text_frame.paragraphs[0].runs[0]
        r.font.size, r.font.bold, r.font.color.rgb = Pt(14), True, BRAND
        tt = slide.shapes.add_textbox(Inches(0.9), Inches(0.72), Inches(11.5), Inches(0.8))
        tt.text_frame.text = spec["title"]
        r = tt.text_frame.paragraphs[0].runs[0]
        r.font.size, r.font.bold, r.font.color.rgb = Pt(32), True, INK
        scale, ox, oy = 0.0104, 3.32, 1.65

        def box(x, y, w, h, fill):
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(ox + x * scale), Inches(oy + y * scale),
                                         Inches(w * scale), Inches(h * scale))
            shp.fill.solid(); shp.fill.fore_color.rgb = fill
            shp.line.fill.background(); shp.shadow.inherit = False
            return shp

        box(44, 14, 552, 412, RGBColor(0x0F, 0x1F, 0x27))
        box(60, 54, 520, 46, RGBColor(0x16, 0x30, 0x3A))
        box(72, 68, 150, 20, RGBColor(0x4A, 0x6B, 0x78))
        box(60, 110, 520, 190, RGBColor(0x0B, 0x17, 0x1D))
        box(72, 122, 300, 42, RGBColor(0x1B, 0x35, 0x40))
        box(268, 176, 300, 42, RGBColor(0x12, 0x3A, 0x32))
        box(72, 230, 330, 42, RGBColor(0x1B, 0x35, 0x40))
        box(70, 352, 410, 34, RGBColor(0x16, 0x30, 0x3A))
        box(490, 352, 82, 34, BRAND)
        box(72, 410, 170, 10, RGBColor(0x2A, 0x45, 0x50))
        x, y, w, h, _rx = _mock_regions().get(spec["region"], _mock_regions()["page"])
        ring = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(ox + (x - 4) * scale), Inches(oy + (y - 4) * scale),
                                      Inches((w + 8) * scale), Inches((h + 8) * scale))
        ring.fill.background(); ring.line.color.rgb = YELLOW; ring.line.width = Pt(3.5)
        ring.shadow.inherit = False
        cb = slide.shapes.add_textbox(Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.6))
        cb.text_frame.word_wrap = True
        cb.text_frame.text = spec["cap"]
        r = cb.text_frame.paragraphs[0].runs[0]
        r.font.size, r.font.color.rgb = Pt(20), RGBColor(0x0A, 0x62, 0x99)

    def code_slide(deck, filename, start, lines, marks, gone, part, parts):
        """A dark, monospaced 'type this now' slide, matching the editor."""
        slide = deck.slides.add_slide(deck.slide_layouts[6])   # blank
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK

        head = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.8))
        label = f"Type this into {filename}"
        if parts > 1:
            label += f"  ({part} of {parts})"
        head.text_frame.text = label
        run = head.text_frame.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.color.rgb = Pt(26), True, BRAND

        box = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12.2), Inches(5.9))
        frame = box.text_frame
        frame.word_wrap = False
        total = len(lines) + sum(len(v) for v in gone.values())
        size = Pt(17) if total <= 11 else Pt(14)
        first = [True]   # the first paragraph already exists; the rest are added

        def row(text, colour, gutter_text, struck=False):
            para = frame.paragraphs[0] if first[0] else frame.add_paragraph()
            first[0] = False
            para.space_after = Pt(0)
            edge = para.add_run()
            edge.text = gutter_text
            edge.font.name, edge.font.size, edge.font.color.rgb = "Consolas", size, GUTTER
            code = para.add_run()
            code.text = text if text.strip() else " "
            code.font.name, code.font.size, code.font.color.rgb = "Consolas", size, colour
            if struck:
                # python-pptx has no strikethrough property, so set the
                # attribute the OOXML run properties element already supports.
                code.font._rPr.set("strike", "sngStrike")

        for i, line in enumerate(lines + [None]):
            number = start + i
            # Anything removed from this position is drawn back in where it
            # used to be. Green alone cannot show a week that deletes a line.
            for dropped in gone.get(number, []):
                row(dropped, DROP, "   -  ", struck=True)
            if line is None:
                break
            # green = new this week, matching the workbook and the week pages
            row(line, NEW if number in marks else PAPER, f"{number:>4}  ")

    def ctx_slide(deck, filename, start, lines, hi, note):
        """The block written so far, with line `hi` highlighted (or every line
        green when hi is None, for the whole-chunk slide), and a note below."""
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = DARK
        head = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.6))
        head.text_frame.text = (f"Type this into {filename}   -   line {start + hi}"
                                if hi is not None else f"All together in {filename}")
        hr = head.text_frame.paragraphs[0].runs[0]
        hr.font.size, hr.font.bold, hr.font.color.rgb = Pt(22), True, BRAND
        box = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.2), Inches(4.7))
        frame = box.text_frame; frame.word_wrap = False
        size = Pt(18) if len(lines) <= 10 else Pt(13)
        first = [True]
        for j, line in enumerate(lines):
            para = frame.paragraphs[0] if first[0] else frame.add_paragraph()
            first[0] = False; para.space_after = Pt(0)
            gut = para.add_run(); gut.text = f"{start + j:>4}  "
            gut.font.name, gut.font.size, gut.font.color.rgb = "Consolas", size, GUTTER
            code = para.add_run(); code.text = line if line.strip() else " "
            code.font.name, code.font.size = "Consolas", size
            # highlighted (or whole-block) line is green; already-typed context is muted
            code.font.color.rgb = GUTTER if (hi is not None and j != hi) else NEW
        nb = slide.shapes.add_textbox(Inches(0.6), Inches(6.15), Inches(12.1), Inches(1.15))
        nf = nb.text_frame; nf.word_wrap = True
        nf.paragraphs[0].text = note
        nr = nf.paragraphs[0].runs[0]; nr.font.size, nr.font.color.rgb = Pt(20), PAPER

    def del_slide(deck, filename, start, context, dropped, note, replacing=False):
        """The block with one line struck through: the line to remove this week."""
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = DARK
        head = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.6))
        head.text_frame.text = f"In {filename}   -   {'replace a line' if replacing else 'delete a line'}"
        hr = head.text_frame.paragraphs[0].runs[0]
        hr.font.size, hr.font.bold, hr.font.color.rgb = Pt(22), True, BRAND
        box = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.2), Inches(4.7))
        frame = box.text_frame; frame.word_wrap = False
        rows = list(context) + [dropped]
        size = Pt(18) if len(rows) <= 10 else Pt(13)
        first = [True]
        for j, line in enumerate(rows):
            para = frame.paragraphs[0] if first[0] else frame.add_paragraph()
            first[0] = False; para.space_after = Pt(0)
            struck = j == len(rows) - 1
            # The struck line sits right after the context, so start + j is its
            # real number - show it, not a dash, so the line is findable.
            gut = para.add_run(); gut.text = f"{start + j:>4}  "
            gut.font.name, gut.font.size, gut.font.color.rgb = "Consolas", size, GUTTER
            code = para.add_run(); code.text = line if line.strip() else " "
            code.font.name, code.font.size, code.font.color.rgb = "Consolas", size, (DROP if struck else GUTTER)
            if struck:
                code.font._rPr.set("strike", "sngStrike")
        nb = slide.shapes.add_textbox(Inches(0.6), Inches(6.15), Inches(12.1), Inches(1.15))
        nf = nb.text_frame; nf.word_wrap = True
        nf.paragraphs[0].text = note
        nr = nf.paragraphs[0].runs[0]; nr.font.size, nr.font.color.rgb = Pt(20), PAPER

    made = 0
    seen = set()
    for week in course.WEEKS:
        deck = Presentation()
        deck.slide_width, deck.slide_height = Inches(13.333), Inches(7.5)
        concept_slide(deck, {"title": f"Week {week['n']}", "sub": week["title"], "bullets": []})
        for kind, d in slide_plan(week, seen):
            if kind == "concept":
                concept_slide(deck, {"title": d["title"], "sub": d["sub"], "bullets": d["bullets"]}, eyebrow=d.get("eyebrow", ""))
            elif kind == "mockup":
                mockup_slide(deck, d)
            elif kind == "checkpoint":
                # No live iframe in PowerPoint, so it points at the classroom.
                concept_slide(deck, {"title": d["title"], "sub": d["say"],
                                     "bullets": ["Run it in the classroom, or press Show the output in the web deck.",
                                                 "Everyone's screen should match this."]},
                              eyebrow="Checkpoint - run it")
            elif kind == "chunk":
                code_slide(deck, d["file"], d["start"], d["lines"], d["marks"], d["gone"], d["part"], d["parts"])
            elif kind == "block":
                ctx_slide(deck, d["file"], d["start"], d["lines"], None,
                          "That is the whole piece. Check yours looks the same before moving on.")
            elif kind == "delete":
                del_slide(deck, d["file"], d["start"], d["lines"], d["dropped"], d["note"],
                          d.get("replacing", False))
            elif kind in ("quiz", "quizanswer"):
                q = d["quiz"]
                answered = kind == "quizanswer"
                letters = "ABCDEFGH"
                bullets = [f"{letters[i]}. {opt}" + (" (correct)" if answered and i == q["answer"] else "")
                           for i, opt in enumerate(q["options"])]
                bullets.append(q["why"] if answered else "Pick one - the answer is on the next slide.")
                concept_slide(deck, {"title": q["q"], "sub": "", "bullets": bullets},
                              eyebrow="Answer" if answered else "Quick check")
            else:  # linectx
                ctx_slide(deck, d["file"], d["start"], d["lines"], d["hi"], d["note"])
        deck.save(os.path.join(SLIDES_DIR, f"week-{week['n']:02d}.pptx"))
        made += 1
    return made


# --------------------------------------------------------------------------
# the same deck as a web page
# --------------------------------------------------------------------------

DECK_CSS = """
*{box-sizing:border-box}
html,body{height:100%;margin:0}
body{background:#0d1b21;color:#1f2a37;font:16px/1.55 Rubik,system-ui,-apple-system,"Segoe UI",sans-serif;
overflow:hidden}
.deck{position:absolute;inset:0}
.slide{position:absolute;inset:0;display:none;flex-direction:column;justify-content:center;
padding:min(6vh,58px) min(6vw,72px) calc(min(6vh,58px) + 54px);background:#f4f8fb}
.slide.on{display:flex}
.slide.dark{background:#102127;color:#dcecef;justify-content:flex-start}
.eyebrow{color:#01aefd;font-size:clamp(13px,1.5vw,19px);font-weight:800;letter-spacing:.09em;
text-transform:uppercase;margin:0 0 .5em}
.slide h1{font-size:clamp(38px,7vw,92px);line-height:1.05;margin:0;font-weight:800;letter-spacing:-.02em}
.slide h2{font-size:clamp(28px,4.4vw,58px);line-height:1.1;margin:0;font-weight:800;letter-spacing:-.015em}
.sub{color:#0a6299;font-size:clamp(19px,2.4vw,33px);font-weight:600;margin:.55em 0 0}
.sub.mono{font-family:Consolas,"SF Mono",Menlo,monospace;font-weight:500;letter-spacing:-.01em}
ul.pts{list-style:none;margin:clamp(20px,3.4vh,44px) 0 0;padding:0;display:grid;
gap:clamp(10px,1.9vh,22px)}
ul.pts li{position:relative;padding-left:1.5em;font-size:clamp(18px,2.3vw,32px);line-height:1.4}
ul.pts li:before{content:"";position:absolute;left:0;top:.52em;width:.62em;height:.62em;
border-radius:50%;background:#01aefd}
.filebar{display:flex;flex-wrap:wrap;gap:10px;align-items:baseline;
font-size:clamp(16px,2vw,27px);font-weight:700;color:#01aefd;margin:0 0 .6em}
.filebar .part{color:#7f9aa4;font-weight:600;font-size:.72em}
.code{font-family:Consolas,"SF Mono",Menlo,monospace;border-collapse:collapse;width:100%;
font-size:clamp(11px,1.55vw,23px);line-height:1.42}
.code td{padding:0;white-space:pre;vertical-align:top}
.code .ln{width:3.4em;text-align:right;padding-right:1.1em;color:#6d8c97;user-select:none}
.code .new td:last-child{color:#7fe0a0}
.code .gone td:last-child{color:#f3928b;text-decoration:line-through}
.code .gone .ln{color:#a8635e}
.bar{position:absolute;left:0;right:0;bottom:0;height:54px;display:flex;align-items:center;
gap:12px;padding:0 18px;background:rgba(16,33,39,.9);color:#dcecef;
font-size:13px;backdrop-filter:blur(6px)}
.bar button{font:inherit;font-weight:600;color:#dcecef;background:#22424d;border:0;
border-radius:6px;padding:7px 13px;cursor:pointer}
.bar button:hover{background:#2d5361}
.bar .count{font-variant-numeric:tabular-nums;letter-spacing:.03em}
.bar .spacer{flex:1}
.bar .hint{color:#8fa9b3}
.dots{display:flex;gap:5px}
.dots i{width:7px;height:7px;border-radius:50%;background:#3a5c68;cursor:pointer}
.dots i.on{background:#01aefd}
.slide.line{overflow-y:auto}
.slide.line .filebar{margin-bottom:.55em}
.slide.line .code{font-size:clamp(13px,1.9vw,26px);margin:0 0 .7em}
.code tr.ctx td{color:#aebfcb}
.code tr.ctx .ln{color:#7c909d}
.code tr.hi td:last-child{color:#7fe0a0;background:rgba(127,224,160,.14);
box-shadow:-.5em 0 0 rgba(127,224,160,.14),4px 0 0 rgba(127,224,160,.14)}
.code tr.hi .ln{color:#7fe0a0;background:rgba(127,224,160,.14);font-weight:700}
.code tr.gone.hi td:last-child{color:#f3928b;background:rgba(243,146,139,.16);text-decoration:line-through}
.code tr.gone.hi .ln{color:#f3928b;background:rgba(243,146,139,.16)}
.linenote{color:#dcecef;font-size:clamp(16px,2vw,28px);line-height:1.45;margin:0;max-width:40ch}
.slide.line.whole .linenote{color:#8fa9b3;font-size:clamp(14px,1.6vw,22px)}
.slide.quiz{justify-content:flex-start;overflow-y:auto}
.slide.quiz .eyebrow{color:#8a5cd0}
.slide.quiz.answered .eyebrow{color:#0f8a4c}
.opts{list-style:upper-alpha;margin:clamp(16px,3vh,34px) 0 0;padding-left:1.5em;
display:grid;gap:clamp(9px,1.7vh,18px)}
.opts li{font-size:clamp(17px,2.2vw,29px);line-height:1.35;padding-left:.3em}
.slide.quiz.answered .opts li{color:#8a9aa5}
.slide.quiz.answered .opts li.right{color:#0f8a4c;font-weight:800}
.quiz-hint{color:#5b7178;font-size:clamp(14px,1.7vw,22px);margin-top:clamp(16px,3vh,30px)}
.quiz-why{color:#0a6299;font-size:clamp(16px,2vw,27px);line-height:1.5;
margin-top:clamp(16px,3vh,30px);max-width:46ch}
.slide.checkpoint{justify-content:flex-start;overflow-y:auto}
.cp-say{color:#0a6299;font-size:clamp(16px,2.1vw,27px);line-height:1.5;margin:.5em 0 1em;max-width:44ch}
.cp-run{font:inherit;font-weight:800;color:#04222f;background:#ffd633;border:0;border-radius:8px;
padding:12px 22px;cursor:pointer;font-size:clamp(15px,1.8vw,22px)}
.cp-run:hover{background:#ffdf5c}
.cp-out{margin-top:16px;width:100%}
.cp-out iframe{width:100%;height:min(42vh,440px);border:1px solid #cbd9dd;border-radius:8px;background:#fff}
.cp-console{margin-top:12px;border:1px solid #21414c;border-radius:8px;overflow:hidden;background:#0f1b21}
.cp-console-head{padding:6px 12px;background:#16303a;color:#cfe3e8;font-size:13px;font-weight:700;
letter-spacing:.03em}
.cp-console-body{max-height:min(22vh,180px);overflow:auto;padding:8px 12px;
font:13px/1.5 Consolas,"SF Mono",Menlo,monospace;color:#dcecef}
.cp-console-body .cp-muted{color:#7f9aa4}
.cp-row{white-space:pre-wrap;overflow-wrap:anywhere;padding:1px 0}
.cp-row.error{color:#ff9d8e}.cp-row.warn{color:#f4c869}.cp-row.net{color:#8fb3bd}.cp-row.info{color:#9fd8ee}
@media (max-width:640px){.bar .hint{display:none}}
.slide.concept-vis{align-items:center;text-align:center;justify-content:center}
.slide.concept-vis .sub{margin-top:.4em}
.vis{width:100%;max-width:min(780px,94%);margin:clamp(12px,2.6vh,26px) auto 0}
.vis.mock{max-width:min(500px,66%)}
.vis svg{width:100%;height:auto;display:block;overflow:visible}
.vis-cap{color:#0a6299;text-align:center;font-size:clamp(16px,2.1vw,27px);font-weight:600;
margin:clamp(8px,1.8vh,18px) auto 0;max-width:40ch}
.vb-rect{fill:#132833;stroke:#33586a;stroke-width:2.5}
.vb-txt{fill:#e6f2f6;font:700 30px Consolas,"SF Mono",Menlo,monospace;text-anchor:middle;dominant-baseline:central}
.vb-in .vb-rect{fill:#20301b;stroke:#ffd633}
.vb-hi{fill:rgba(1,174,253,.30);stroke:#01aefd;stroke-width:3}
.vmachine{fill:#0f2c38;stroke:#01aefd;stroke-width:3}
.vmlabel{fill:#7fd8ff;font:800 26px Rubik,system-ui,sans-serif;text-anchor:middle;dominant-baseline:central}
.vtile rect{fill:#16303a;stroke:#33586a;stroke-width:2.5}
.vtile text{fill:#e6f2f6;font:700 23px Consolas,"SF Mono",Menlo,monospace;text-anchor:middle;dominant-baseline:central}
.vt-in rect{stroke:#ffd633}
.vt-out rect{fill:#16321f;stroke:#7fe0a0}
.vt-cloud rect{fill:#0f2c38;stroke:#01aefd}
.vplus{fill:#7fd8ff;font:800 42px Rubik,system-ui,sans-serif;text-anchor:middle;dominant-baseline:central}
.varrow{fill:#5b8aa0;font:800 30px Rubik,system-ui,sans-serif;text-anchor:middle;dominant-baseline:central}
.vswap-off{fill:#2a2030;stroke:#c98be0;stroke-width:3}
.vswap-on{fill:#16321f;stroke:#7fe0a0;stroke-width:3}
.vb-pick .vb-rect{fill:#16321f;stroke:#7fe0a0}
.vt-btn rect{fill:#0f2c38;stroke:#01aefd}
.vt-tag rect{fill:#0f2c38;stroke:#01aefd}
.vt-tag text{fill:#7fd8ff}
.vt-no rect{fill:#1a2730;stroke:#7c909d}
.vt-no text{fill:#aebfcb}
"""

DECK_JS = r"""
var slides = [].slice.call(document.querySelectorAll('.slide'));
var dots = [].slice.call(document.querySelectorAll('.dots i'));
var count = document.querySelector('.count');
var at = 0;

function show(next, push) {
  at = Math.max(0, Math.min(slides.length - 1, next));
  slides.forEach(function (s, i) { s.classList.toggle('on', i === at); });
  dots.forEach(function (d, i) { d.classList.toggle('on', i === at); });
  count.textContent = (at + 1) + ' / ' + slides.length;
  if (push !== false) history.replaceState(null, '', '#' + (at + 1));
  focusTypedLine(slides[at]);
}

// A code slide can be taller than the screen. The line to type is always the
// newest (last) line of the block, so when a scrollable code slide opens, jump
// to the bottom - the line they type and its note land on screen and nobody has
// to scroll down to find them. Review ("all together") slides have no such line,
// so they read from the top.
function focusTypedLine(slide) {
  if (!slide || !slide.classList.contains('line')) return;
  var hasTarget = !!slide.querySelector('tr.hi');
  function go() {
    // Bail if the class moved on before a late callback (font load) fires.
    if (slides[at] !== slide || !slide.classList.contains('on')) return;
    slide.scrollTop = hasTarget ? slide.scrollHeight : 0;
  }
  requestAnimationFrame(go);
  // Rubik/Consolas can load after first paint and make the code taller; without
  // re-running once fonts are in, the jump uses the pre-font height and stops
  // short, leaving the line off-screen - exactly the bug this fixes.
  if (document.fonts && document.fonts.ready) document.fonts.ready.then(go);
}
// Deep link, so a reload - or opening straight at a slide - keeps the place.
show(parseInt((location.hash || '').slice(1), 10) - 1 || 0, false);

document.querySelector('.prev').onclick = function () { show(at - 1); };
document.querySelector('.next').onclick = function () { show(at + 1); };
dots.forEach(function (d, i) { d.onclick = function () { show(i); }; });
document.querySelector('.full').onclick = function () {
  if (document.fullscreenElement) document.exitFullscreen();
  else document.documentElement.requestFullscreen();
};
// Full screen and any projector/window resize change the slide height, so
// re-pin the line-to-type to the bottom of whatever code slide is showing.
window.addEventListener('resize', function () { focusTypedLine(slides[at]); });

document.addEventListener('keydown', function (event) {
  var key = event.key;
  if (key === 'ArrowRight' || key === 'PageDown' || key === ' ' || key === 'Enter') show(at + 1);
  else if (key === 'ArrowLeft' || key === 'PageUp' || key === 'Backspace') show(at - 1);
  else if (key === 'Home') show(0);
  else if (key === 'End') show(slides.length - 1);
  else if (key === 'f') document.querySelector('.full').click();
  else if (key === 'Escape') {
    // Inside the classroom this deck is an iframe: tell the page holding it to
    // close rather than leaving the teacher stuck in a panel with no exit.
    if (document.fullscreenElement) return;   // the browser handles that one
    if (window.parent !== window) window.parent.postMessage({ utg: 'deck', action: 'close' }, location.origin);
  } else return;
  event.preventDefault();
});
// An iframe gets no keys until something in it is focused.
window.addEventListener('load', function () { window.focus(); });
document.addEventListener('click', function () { window.focus(); });

// The demo AI key, so a checkpoint on an AI week can reply for real. It lives in
// the classroom database, never in these published files, and is handed only to
// a TEACHER. A teacher proves it either with a signed-in instructor/admin
// account token, or - the common case - with the instructor 4-letter code they
// used to open the curriculum, which the access gate saved as utg_class_code. A
// student has a STUDENT code, which the server refuses, so the checkpoint just
// shows the interface for them. Fetched once and cached.
var _cpKey;
function cpDemoKey() {
  if (_cpKey) return _cpKey;
  _cpKey = new Promise(function (resolve) {
    var host = location.hostname;
    var api = (host === 'localhost' || host === '127.0.0.1')
      ? 'http://127.0.0.1:8787' : 'https://utg-classroom-api.utgapps.workers.dev';
    var token = null, code = null;
    try { token = (JSON.parse(localStorage.getItem('utg_account') || 'null') || {}).token; } catch (e) {}
    try { code = (localStorage.getItem('utg_class_code') || '').trim(); } catch (e) {}
    var req;
    if (token) {
      req = fetch(api + '/demo-key', { cache: 'no-store', headers: { authorization: 'Bearer ' + token } });
    } else if (code) {
      req = fetch(api + '/demo-key', { method: 'POST', cache: 'no-store',
        headers: { 'content-type': 'application/json' }, body: JSON.stringify({ code: code }) });
    } else { resolve(null); return; }
    req.then(function (r) { return r.ok ? r.json() : null; })
      .then(function (d) { resolve((d && d.key) || null); })
      .catch(function () { resolve(null); });
  });
  return _cpKey;
}

// Checkpoint slides: reveal the project's live output so the class can compare.
[].slice.call(document.querySelectorAll('.cp-run')).forEach(function (btn) {
  btn.onclick = function (e) {
    e.stopPropagation();
    var id = btn.getAttribute('data-cp');
    var out = document.getElementById('out-' + id);
    if (out.firstChild) {
      if (out._cpStop) out._cpStop();       // drop the console listener
      out.innerHTML = ''; btn.textContent = '▶ Show the output'; return;
    }
    // The stored page escapes the script close tag so the text/plain box is not
    // closed early; undo that before it becomes a real document, or the first
    // script tag never closes and none of the project's JS runs. (The close tag
    // is built from pieces so this very file does not contain a literal one.)
    var CLOSE = '<' + '/script';
    var raw = document.getElementById('src-' + id).textContent.split('<\\' + '/script').join(CLOSE);
    var frame = document.createElement('iframe');
    frame.setAttribute('sandbox', 'allow-scripts allow-forms');
    frame.setAttribute('allow', 'local-network-access *');
    out.appendChild(frame);

    // A console under the output, fed by the bridge inside the page. It is how
    // the class sees a console.log, an error, or a request that failed.
    var con = document.createElement('div');
    con.className = 'cp-console';
    con.innerHTML = '<div class="cp-console-head">Console</div>'
                  + '<div class="cp-console-body"><span class="cp-muted">console.log, errors '
                  + 'and network requests from the output show up here.</span></div>';
    out.appendChild(con);
    var body = con.querySelector('.cp-console-body');
    var first = true;
    function onMsg(ev) {
      if (ev.source !== frame.contentWindow) return;          // only this output
      var d = ev.data; if (!d || d.utg !== 'cp-log') return;
      if (first) { body.innerHTML = ''; first = false; }       // clear the hint
      var row = document.createElement('div');
      row.className = 'cp-row ' + (d.level || 'log');
      row.textContent = d.text;
      body.appendChild(row); body.scrollTop = body.scrollHeight;
    }
    window.addEventListener('message', onMsg);
    out._cpStop = function () { window.removeEventListener('message', onMsg); };

    // If this page calls the AI, put the teacher's demo key where the student's
    // key would go, so it can reply for real. No key (or not a teacher) leaves
    // the placeholder - the interface still runs, the request just says why.
    var needsKey = raw.indexOf('sk-class-put-your-own-key-here') >= 0;
    (needsKey ? cpDemoKey() : Promise.resolve(null)).then(function (key) {
      frame.srcdoc = key ? raw.split('sk-class-put-your-own-key-here').join(key) : raw;
    });
    btn.textContent = 'Hide the output';
  };
});
"""


def deck_code_slide(filename, start, lines, marks, gone, part, parts):
    """One dark 'type this now' slide, matching the editor and the pptx."""
    rows = []
    for index, line in enumerate(lines + [None]):
        number = start + index
        for dropped in gone.get(number, []):
            rows.append('<tr class="gone"><td class="ln">{0}</td><td>{1}</td></tr>'
                        .format(number, esc(dropped) or "&nbsp;"))
        if line is None:
            break
        rows.append('<tr class="{0}"><td class="ln">{1}</td><td>{2}</td></tr>'.format(
            "new" if number in marks else "", number, esc(line) or "&nbsp;"))
    part_label = ('<span class="part">{0} of {1}</span>'.format(part, parts)) if parts > 1 else ""
    return ('<section class="slide dark">'
            '<p class="filebar">Type this into {0} {1}</p>'
            '<table class="code">{2}</table></section>').format(esc(filename), part_label, "".join(rows))


def deck_concept_slide(spec):
    sub = ""
    if spec.get("sub"):
        # A sub-heading that is really a line of code should look like one.
        code_ish = any(mark in spec["sub"] for mark in ("(", "{", "[", "=", ".", "/", ":", ";", "<"))
        sub = '<p class="sub{0}">{1}</p>'.format(" mono" if code_ish else "", esc(spec["sub"]))
    points = "".join("<li>{0}</li>".format(esc(point)) for point in spec.get("bullets", []))
    return ('<section class="slide"><h2>{0}</h2>{1}{2}</section>'.format(
        esc(spec["title"]), sub,
        '<ul class="pts">{0}</ul>'.format(points) if points else ""))


def build_html_decks():
    """The same deck as a web page, so it can be presented from the classroom.

    Built from the same week["slides"] and the same code_chunks() as the .pptx,
    slide for slide, and main() fails the build if the two ever disagree. The
    .pptx stays for Google Slides; this is what the teacher screen embeds.
    """
    os.makedirs(SLIDES_DIR, exist_ok=True)
    counts = {}
    seen = set()
    for week in course.WEEKS:
        slides = ['<section class="slide"><p class="eyebrow">Week {0}</p><h1>{1}</h1></section>'
                  .format(week["n"], esc(week["title"]))]
        for desc in slide_plan(week, seen):
            slides.append(deck_render_html(desc))
        dots = "".join("<i></i>" for _ in slides)
        html_out = """<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Week {n} slides &middot; AI101 &middot; UTG Academy</title>
{font}
<style>{css}</style>
{guard}
<div class="deck">{slides}</div>
<div class="bar">
  <button class="prev" title="Previous slide (left arrow)">&lsaquo;</button>
  <button class="next" title="Next slide (right arrow, space)">&rsaquo;</button>
  <span class="count"></span>
  <div class="dots">{dots}</div>
  <span class="spacer"></span>
  <span class="hint">arrows or space to move &middot; f for full screen</span>
  <button class="full" title="Full screen (f)">Full screen</button>
</div>
<script>{js}</script>
""".format(n=week["n"], font=FONT, css=DECK_CSS, guard=guard("ai101", up="../../"),
           slides="".join(slides), dots=dots, js=DECK_JS)
        with open(os.path.join(SLIDES_DIR, "week-%02d.html" % week["n"]), "w",
                  encoding="utf-8") as handle:
            handle.write(html_out)
        counts[week["n"]] = len(slides)
    return counts


def pptx_slide_counts():
    """Slide count per week, from the same plan the decks render, so the
    .pptx and .html can never disagree."""
    counts = {}
    seen = set()
    for week in course.WEEKS:
        counts[week["n"]] = 1 + len(slide_plan(week, seen))
    return counts

def build_milestones():
    """Emit every week's file set as JSON, for the classroom editor to read.

    This is what lets an instructor catch up a student who missed a week: the
    editor fetches this and seeds week N's canonical code into their account.
    Public on purpose - it is the same code already printed on the week pages,
    and it contains no key, only the placeholder every student replaces.
    """
    payload = {
        "course": "ai101",
        "title": course.COURSE_TITLE,
        "weeks": [
            {"n": week["n"], "title": week["title"], "files": state_at(week["n"])}
            for week in course.WEEKS
        ],
    }
    for week in payload["weeks"]:
        for name, text in week["files"].items():
            if "sk-class-" in text and "put-your-own-key-here" not in text:
                raise SystemExit(f"week {week['n']} {name} carries a real-looking key; refusing to publish")
    write("milestones.json", json.dumps(payload, separators=(",", ":")))
    return len(payload["weeks"])


def write(name, text):
    with open(os.path.join(HERE, name), "w", encoding="utf-8") as handle:
        handle.write(text)


# --------------------------------------------------------------------------

def main():
    if len(course.WEEKS) != 15:
        raise SystemExit(f"expected 15 weeks, found {len(course.WEEKS)}")
    for i, week in enumerate(course.WEEKS, start=1):
        if week["n"] != i:
            raise SystemExit(f"weeks are out of order at position {i}")

    final = state_at(15)
    total = line_count(final)
    print(f"AI101 - final project is {total} lines "
          f"({', '.join(f'{n}: {len(final[n].splitlines())}' for n in FILES)})")
    if total > LINE_BUDGET:
        raise SystemExit(f"OVER BUDGET: {total} lines > {LINE_BUDGET}. Move something to a bonus module.")

    # A week that adds nothing is almost always an authoring slip.
    for week in course.WEEKS:
        if not week["ops"] and not week.get("no_code_ok"):
            raise SystemExit(f"week {week['n']} changes no code; set no_code_ok=True if that is deliberate")

        # Every edit must be pinned to a moment in the lesson, or the teacher is
        # left guessing when in the hour it happens.
        touched = {(filename, block_id) for _, filename, block_id, _ in week["ops"]}
        cited = {(beat["file"], beat["block"]) for beat in week["flow"] if beat["kind"] == "step"}
        orphans = touched - cited
        if orphans:
            raise SystemExit(
                f"week {week['n']}: no STEP in the flow types "
                + ", ".join(f"{f}:{b}" for f, b in sorted(orphans))
            )
        ghosts = {ref for ref in cited if ref[0] not in FILES}
        if ghosts:
            raise SystemExit(f"week {week['n']}: flow cites unknown file(s) {sorted(ghosts)}")

        # A whole hour with nothing to say back is a lecture, not a lesson.
        prompts = sum(1 for beat in week["flow"] if beat.get("ask"))
        if prompts < 3:
            raise SystemExit(
                f"week {week['n']}: only {prompts} call-and-response prompt(s) in the flow; "
                f"aim for at least 3 spread through the hour"
            )

        # The deck is what is on the projector while students type, so every
        # edit has to be on a slide too - not only in the teacher's own notes.
        on_slides = {ref for spec in week["slides"] for ref in spec.get("code", [])}
        unshown = touched - on_slides
        if unshown:
            raise SystemExit(
                f"week {week['n']}: no slide shows "
                + ", ".join(f"{f}:{b}" for f, b in sorted(unshown))
                + " - add a \"code\" key to one of that week's slides"
            )

        # An expanded deck IS the type-along, so its code must appear in the
        # order the lesson types it - not the author's concept order. Otherwise
        # a slide says "type line 24" before "line 9".
        if week["n"] in getattr(course, "EXPANDED_WEEKS", set()):
            slide_blocks = [ref for spec in week["slides"] for ref in spec.get("code", [])]
            flow_blocks = [(beat["file"], beat["block"]) for beat in week["flow"]
                           if beat["kind"] == "step"]
            if slide_blocks != flow_blocks:
                raise SystemExit(
                    f"week {week['n']} is expanded but its slide code order "
                    f"{slide_blocks} does not match the typing order {flow_blocks} - "
                    f"reorder that week's slides so the type-along follows the lesson"
                )

        # An expanded week teaches line by line, so every code line it shows on
        # a slide needs a note. A blank explanation slide is worse than none, so
        # refuse to ship one - this is the guard that keeps the roll-out honest.
        if week["n"] in getattr(course, "EXPANDED_WEEKS", set()):
            for spec in week["slides"]:
                for filename, block in (spec.get("code") or []):
                    start, lines, marks, _gone = block_code(week["n"], filename, block)
                    notes = notes_for(week["n"], filename, block)
                    for i, line in enumerate(lines):
                        # Only a changed/new line gets a slide, so only it needs a
                        # note; a carried-over line just shows as context.
                        if not line.strip() or (start + i) not in marks:
                            continue
                        note = notes[i] if notes and i < len(notes) else ""
                        if note is None:
                            continue
                        if not note:
                            raise SystemExit(
                                f"week {week['n']} is expanded but {filename}:{block} "
                                f"line {start + i} has no LINE_NOTES entry"
                            )

        grew = line_count(state_at(week["n"])) - (line_count(state_at(week["n"] - 1)) if week["n"] > 1 else 0)
        print(f"  week {week['n']:2d}  +{grew:3d} lines  {week['title']}")

    # A quiz with a bad answer index marks no option correct, silently - refuse.
    for key, quiz_set in getattr(course, "QUIZZES", {}).items():
        for q in quiz_set:
            opts = q.get("options") or []
            if not q.get("q") or len(opts) < 2 or not q.get("why"):
                raise SystemExit(f"quiz {key}: needs q, why, and 2+ options - {q.get('q')!r}")
            if not isinstance(q.get("answer"), int) or not (0 <= q["answer"] < len(opts)):
                raise SystemExit(f"quiz {key}: answer index out of range - {q.get('q')!r}")

    build_index()
    build_weeks()
    build_teacher()
    build_workbook()
    weeks_out = build_milestones()
    decks = build_slides()
    web = build_html_decks()
    # The .pptx is for Google Slides and the .html is what the classroom
    # embeds. They are generated from the same week data, so if they ever
    # disagree on slide count, one of the two renderers has drifted.
    expected = pptx_slide_counts()
    drift = {n: (web[n], expected[n]) for n in expected if web[n] != expected[n]}
    if drift:
        raise SystemExit("deck drift - week: (html, pptx) " + repr(drift))
    print(f"built: index.html, week-01..15.html, teacher.html, workbook.html, "
          f"{decks} slide decks (.pptx and .html, {sum(web.values())} slides), "
          f"milestones.json ({weeks_out} weeks)")


if __name__ == "__main__":
    main()
